"""
Template Agent
负责解析 PDF/PPTX 模版，提取风格、Logo 和关键页面参考图
"""
import os
import json
import fitz  # PyMuPDF
from PIL import Image
import logging
from typing import Dict, List, Optional
from openai import OpenAI

from ..utils.llm_client import chat_completion_with_fallback, MODEL_FALLBACK_CHAIN

import base64
import io

# Try to import pptx for pptx parsing support
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TemplateAgent:
    def __init__(self, api_key: str, api_base: str = None, output_dir: str = "template_assets"):
        self.client = OpenAI(
            api_key=api_key,
            base_url=api_base or "https://generativelanguage.googleapis.com/v1beta/openai",
            timeout=120.0,
            max_retries=3
        )
        self.model = "gemini-3.1-pro-preview"
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def process_template(self, template_path: str) -> Dict:
        """
        处理模版文件，返回模版资产信息
        """
        logger.info(f"🎨 Template Agent: 正在解析模版 {template_path}...")
        
        # Branch based on file extension
        if template_path.lower().endswith('.pptx'):
            return self.process_pptx_template(template_path)
        
        # 1. 将模版转换为图片序列 (PDF fallback)
        images = self._convert_to_images(template_path)
        if not images:
            raise ValueError("无法解析模版文件")
            
        # 2. 使用 AI 分析页面类型并提取风格
        template_info = self._analyze_template_structure(images)
        
        # 3. 提取 Logo (尝试从封面提取)
        logo_path = self._extract_logo(images[0], template_info.get('logo_location'))
        template_info['logo_path'] = logo_path
        
        # 4. 保存关键参考图
        ref_paths = self._save_reference_images(images, template_info.get('page_types', []))
        template_info['reference_images'] = ref_paths
        
        return template_info

    def process_pptx_template(self, pptx_path: str) -> Dict:
        """
        解析 .pptx 模版，提取母版信息和布局
        Also attempts to convert PPTX to PDF using 'soffice' to extract visual assets.
        """
        if not HAS_PPTX:
            logger.error("Missing python-pptx library.")
            raise ImportError("Please install python-pptx")
            
        logger.info("🧠 正在解析 .pptx 母版结构...")
        prs = Presentation(pptx_path)
        
        # Extract basic info
        template_info = {
            "source_type": "pptx",
            "file_path": pptx_path,
            "layouts": [],
            "color_palette": [],
            "fonts": [],
            "page_types": [] # We will infer this
        }
        
        # Analyze Layouts
        # Common layout names mapping to types
        layout_mapping = {
            "Title Slide": "Cover",
            "Title and Content": "Content",
            "Section Header": "Section",
            "Two Content": "Content",
            "Comparison": "Content",
            "Title Only": "Hero",
            "Blank": "Back"
        }
        
        # Extract Layouts
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = layout.name
            mapped_type = layout_mapping.get(layout_name, "Content") # Default to Content
            
            # Simple heuristic for unmapped names
            if "title" in layout_name.lower() and "content" not in layout_name.lower():
                mapped_type = "Cover"
            elif "content" in layout_name.lower():
                mapped_type = "Content"
                
            template_info["layouts"].append({
                "index": i,
                "name": layout_name,
                "type": mapped_type
            })
            
        # Extract Colors (Simple Heuristic: Check Title Placeholder Color on Master)
        try:
            # Try to find a title style in the first master
            if prs.slide_masters:
                master = prs.slide_masters[0]
                # Check Title Style
                if master.slide_layouts:
                    title_layout = master.slide_layouts[0]
                    if title_layout.placeholders:
                        title_ph = title_layout.placeholders[0]
                        if hasattr(title_ph, 'text_frame'):
                            # Try to get font
                            try:
                                font = title_ph.text_frame.paragraphs[0].font
                                if font.name:
                                    template_info["fonts"].append(font.name)
                                if hasattr(font.color, 'rgb') and font.color.rgb:
                                    template_info["color_palette"].append(f"#{font.color.rgb}")
                            except:
                                pass
        except Exception as e:
            logger.warning(f"Could not extract font/color details from PPTX: {e}")
            
        # Try to convert to PDF for Visual Asset Extraction
        try:
            pdf_path = self._convert_pptx_to_pdf(pptx_path)
            if pdf_path and os.path.exists(pdf_path):
                logger.info(f"✅ PPTX 已转换为 PDF: {pdf_path}. 开始提取视觉资产...")
                
                # Use PDF visual extraction logic
                images = self._convert_to_images(pdf_path)
                
                if images:
                    # Analyze structure using Vision
                    visual_info = self._analyze_template_structure(images)
                    
                    # Merge info
                    template_info["page_types"] = visual_info.get("page_types", [])
                    template_info["logo_location"] = visual_info.get("logo_location")
                    
                    if not template_info["color_palette"]:
                        template_info["color_palette"] = visual_info.get("color_palette", [])
                    
                    # Extract Logo
                    logo_path = self._extract_logo(images[0], template_info.get('logo_location'))
                    template_info['logo_path'] = logo_path
                    
                    # Save Ref Images
                    ref_paths = self._save_reference_images(images, template_info.get('page_types', []))
                    template_info['reference_images'] = ref_paths
                    
                    # Use Vision-based style description as it's richer
                    template_info["style_description"] = visual_info.get("style_description")
                    
                    # Cleanup PDF
                    try:
                        os.remove(pdf_path)
                    except:
                        pass
                        
                    return template_info
        except Exception as e:
            logger.warning(f"PPTX to PDF conversion failed: {e}. Fallback to metadata extraction.")

        # Fallback to AI Analysis for Style Description if we can't render
        # Build a textual style description from extracted data
        palette_str = ", ".join(template_info.get("color_palette", []))
        fonts_str = ", ".join(template_info.get("fonts", []))
        
        template_info["style_description"] = f"""
        Corporate Template (Extracted from PPTX).
        Theme Colors: {palette_str}.
        Fonts: {fonts_str}.
        Layout: Follow the master slide structure strictly.
        """
        
        return template_info

    def _convert_pptx_to_pdf(self, pptx_path: str) -> Optional[str]:
        """
        Convert PPTX to PDF using LibreOffice (soffice)
        Returns path to generated PDF or None
        """
        import subprocess
        
        # Check for soffice
        soffice_path = None
        possible_paths = [
            "/opt/homebrew/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "soffice"
        ]
        
        for p in possible_paths:
            if os.path.exists(p) or (p == "soffice" and os.system("which soffice > /dev/null 2>&1") == 0):
                soffice_path = p
                break
                
        if not soffice_path:
            logger.warning("LibreOffice (soffice) not found. Cannot convert PPTX to PDF.")
            return None
            
        output_dir = os.path.dirname(pptx_path)
        # Handle filename with spaces
        cmd = [
            soffice_path,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            pptx_path
        ]
        
        logger.info(f"Running conversion: {' '.join(cmd)}")
        try:
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode == 0:
                # Assuming output filename is same as input but .pdf
                base_name = os.path.splitext(os.path.basename(pptx_path))[0]
                pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
                if os.path.exists(pdf_path):
                    return pdf_path
            else:
                logger.error(f"soffice conversion failed: {result.stderr}")
        except Exception as e:
            logger.error(f"Subprocess failed: {e}")
            
        return None

    def _convert_to_images(self, path: str) -> List[Image.Image]:
        """将 PDF/PPT 转为 PIL Image 列表"""
        images = []
        if path.lower().endswith('.pdf'):
            doc = fitz.open(path)
            for page in doc:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_data = pix.tobytes("png")
                images.append(Image.open(io.BytesIO(img_data)))
        return images

    def _analyze_template_structure(self, images: List[Image.Image]) -> Dict:
        """调用 Vision 模型分析模版结构"""
        logger.info("🧠 正在分析模版结构与风格...")
        
        # 选取前 5 页进行分析
        sample_images = images[:5]
        
        # 构建多模态 Prompt
        content = [
            {"type": "text", "text": """你是专业的 UI/UX 设计师。请分析这组 PPT 模版图片。
            
            任务 1: 识别每张图的页面类型 (Cover, TOC, Section, Content, Hero, Back)。
            任务 2: 提取全局视觉风格 (配色、字体、装饰)。请特别注意提取准确的 HEX 颜色代码。
            任务 3: 找出 Logo 的位置 (Top-Left, Top-Right, Center 等)。
            
            请输出 JSON 格式：
            {
              "page_types": ["Cover", "Content", "Content", ...],
              "style_description": "Deep blue background...",
              "logo_location": "Top-Right",
              "color_palette": ["#000000", "#FFFFFF"],
              "fonts": ["Arial", "Helvetica"]
            }"""}
        ]
        
        for img in sample_images:
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            b64 = base64.b64encode(buffered.getvalue()).decode()
            content.append({
                "type": "image_url", 
                "image_url": {"url": f"data:image/png;base64,{b64}"}
            })
            
        try:
            response = chat_completion_with_fallback(
                self.client, model=self.model, model_fallback=MODEL_FALLBACK_CHAIN,
                messages=[{"role": "user", "content": content}],
                temperature=0.1
            )
            result = response.choices[0].message.content.strip()
            # Clean JSON
            if result.startswith('```json'): result = result[7:]
            if result.endswith('```'): result = result[:-3]
            return json.loads(result)
        except Exception as e:
            logger.error(f"模版分析失败: {e}")
            return {
                "page_types": ["Cover", "Content", "Content"],
                "style_description": "Professional business style",
                "logo_location": "Top-Right",
                "color_palette": ["#000000", "#FFFFFF"]
            }

    def _extract_logo(self, cover_image: Image.Image, location: str) -> Optional[str]:
        """
        简单的 Logo 提取逻辑 (裁剪)
        实际生产中应用 Object Detection 模型
        """
        if not location:
            return None
            
        w, h = cover_image.size
        crop_box = None
        
        if "Left" in location and "Top" in location:
            crop_box = (0, 0, w//3, h//4)
        elif "Right" in location and "Top" in location:
            crop_box = (w*2//3, 0, w, h//4)
        elif "Center" in location:
            crop_box = (w//3, h//3, w*2//3, h*2//3)
            
        if crop_box:
            try:
                logo_img = cover_image.crop(crop_box)
                save_path = os.path.join(self.output_dir, "extracted_logo.png")
                logo_img.save(save_path)
                logger.info(f"Logo 已提取: {save_path}")
                return save_path
            except Exception as e:
                logger.warning(f"Logo 提取失败: {e}")
                
        return None

    def _save_reference_images(self, images: List[Image.Image], page_types: List[str]) -> Dict:
        """保存分类后的参考图"""
        refs = {}
        # 确保 page_types 长度与 images 一致，不够则补 Content
        if len(page_types) < len(images):
            page_types.extend(['Content'] * (len(images) - len(page_types)))
            
        for img, p_type in zip(images, page_types):
            key = f"ref_{p_type.lower()}"
            if key not in refs: # 只存第一张遇到的该类型
                path = os.path.join(self.output_dir, f"{key}.png")
                img.save(path)
                refs[key] = path
        return refs

if __name__ == "__main__":
    pass
