"""
PPT 生成工作流 - 基于 Banana Slides 思路的简化版
支持从想法/大纲生成完整 PPT，并导出为 PPTX 格式
"""
import os
import json
import base64
import logging
import requests
from pathlib import Path
from typing import List, Dict, Optional, Union
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from openai import OpenAI

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def _fix_black_corners(img: Image.Image, corner_ratio: float = 0.18, dark_threshold: int = 22) -> Image.Image:
    """
    已停用：原逻辑用 np.fliplr 镜像填充右上角黑块，会错误地把左侧文字镜像到右上角，
    导致 P7/P12/P14 等多页出现重复且镜像的文字。改为直接返回原图，依赖 visual_prompt
    的 "NO empty black corners" 约束来避免黑块。
    """
    return img


class PPTGenerator:
    """PPT 生成器 - 核心工作流"""
    
    def __init__(self, api_key: str, api_base: str = None, slides_dir: str = "output/slides"):
        """
        Args:
            api_key: API 密钥
            api_base: API Base URL（可选）
            slides_dir: 临时幻灯片图片的保存目录
        """
        self._api_key = api_key
        self.client = OpenAI(
            api_key=api_key,
            base_url=api_base or "https://generativelanguage.googleapis.com/v1beta/openai",
            timeout=120.0,
            max_retries=3
        )
        self.text_model = "gemini-3-flash-preview"
        self.visual_director_model = "gemini-3.1-pro-preview"
        self.image_model = "gemini-3.1-flash-image-preview"
        self.output_dir = Path(slides_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def generate_image(self, description: str, aspect_ratio: str = "16:9", reference_images: List[Image.Image] = None, is_background_only: bool = False, resolution: str = "1K", native_images: List[Dict] = None) -> Image.Image:
        """
        生成单页PPT图片
        使用 Google REST API 直接调用（兼容 nano banana）
        resolution: "1K" | "2K" | "4K"，默认 1K
        """
        import requests

        resolution = (resolution or "1K").upper()
        if resolution not in ("1K", "2K", "4K"):
            resolution = "1K"

        logger.info(f"正在生成图片 ({resolution}): {description[:50]}...")
        if reference_images:
            logger.info(f"🎨 使用 {len(reference_images)} 张参考图片保持风格一致")

        # VisualAgent 已构造完整 prompt，此处仅追加分辨率等技术参数
        tech_suffix = f"\n\nTechnical: aspect ratio {aspect_ratio}, {resolution} resolution, sharp text rendering. CRITICAL: No black blocks, no solid black rectangles, seamless full-bleed composition."
        
        # Inject smart whitespace instructions based on native_images array
        if native_images and len(native_images) > 0:
            areas = []
            for idx, img_conf in enumerate(native_images):
                layout = img_conf.get('layout')
                bbox = img_conf.get('bounding_box')
                
                if bbox:
                    # Translate bounding box to natural language roughly
                    left_pct = int(bbox.get('left', 0) * 100)
                    top_pct = int(bbox.get('top', 0) * 100)
                    w_pct = int(bbox.get('width', 0) * 100)
                    h_pct = int(bbox.get('height', 0) * 100)
                    
                    # Provide an even stronger spatial instruction
                    if left_pct > 50:
                        position = "the RIGHT SIDE"
                    elif left_pct + w_pct < 50:
                        position = "the LEFT SIDE"
                    else:
                        position = "the CENTER"
                        
                    areas.append(f"a massive empty space on {position} (starting {left_pct}% from left, {top_pct}% from top, spanning {w_pct}% width)")
                elif layout:
                    # Legacy fallback
                    layout_prompts = {
                        "right_half": "the RIGHT SIDE of the image",
                        "left_half": "the LEFT SIDE of the image",
                        "center": "the CENTER area of the image",
                        "bottom_right": "the BOTTOM RIGHT corner"
                    }
                    if layout in layout_prompts:
                        areas.append(layout_prompts[layout])
            
            if areas:
                areas_str = " and ".join(areas)
                tech_suffix += f" CRITICAL VISUAL CONSTRAINT: You ABSOLUTELY MUST leave {areas_str} completely BLANK and EMPTY. Do NOT generate ANY text, shapes, or complex backgrounds in this area. It must be a flat, solid color gradient because a photo will be pasted over it later."



        full_prompt = description + tech_suffix

        # 使用 Google REST API 直接调用
        api_key = self.client.api_key if hasattr(self.client, 'api_key') else None
        if not api_key:
            api_key = getattr(self, '_api_key', None)

        if not api_key:
            raise ValueError("无法获取 API Key")

        url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.image_model}:generateContent?key={api_key}"
        headers = {"Content-Type": "application/json"}

        # 构建 parts 数组：先文本，后参考图片
        parts = [{"text": full_prompt}]

        # 添加参考图片（如果提供）- 用于风格一致性
        if reference_images:
            for ref_img in reference_images[:1]:  # 限制为最多1张，避免干扰文字
                buffered = BytesIO()
                ref_img.save(buffered, format="PNG")
                img_b64 = base64.b64encode(buffered.getvalue()).decode()
                parts.append({
                    "inlineData": {
                        "mimeType": "image/png",
                        "data": img_b64
                    }
                })

        # generationConfig: 支持 1K/2K/4K 分辨率
        generation_config = {
            "responseModalities": ["IMAGE"],
            "imageConfig": {
                "aspectRatio": aspect_ratio,
                "imageSize": resolution
            }
        }
        payload = {
            "contents": [{"parts": parts}],
            "generationConfig": generation_config,
        }
        
        # 添加重试机制（增加延迟，避免限流）
        max_retries = 5  # Increased from 3 to 5
        import time
        
        # 在请求前添加短暂延迟，避免触发 API 限流
        if hasattr(self, '_last_request_time'):
            elapsed = time.time() - self._last_request_time
            if elapsed < 3:  # Increased from 2 to 3 seconds
                time.sleep(3 - elapsed)
        
        for attempt in range(max_retries):
            try:
                response = requests.post(url, headers=headers, json=payload, timeout=180)
                
                if response.status_code == 200:
                    result = response.json()
                    candidates = result.get('candidates', [])
                    
                    if not candidates:
                        # Log the full result for debugging
                        logger.warning(f"Attempt {attempt + 1}: Empty candidates in response: {str(result)[:200]}...")
                        raise ValueError("生成失败：返回结果为空")
                    
                    content = candidates[0].get('content', {})
                    parts = content.get('parts', [])
                    
                    for part in parts:
                        # 兼容 REST API 可能返回的两种字段格式
                        img_b64 = None
                        if 'inlineData' in part:
                            img_b64 = part['inlineData']['data']
                        elif 'inline_data' in part:
                            img_b64 = part['inline_data']['data']
                        
                        if img_b64:
                            img_data = base64.b64decode(img_b64)
                            image = Image.open(BytesIO(img_data)).convert("RGB")
                            image = _fix_black_corners(image)
                            logger.info(f"图片生成成功: {image.size}")
                            self._last_request_time = time.time()
                            return image
                    
                    raise ValueError("未找到图片数据")
                else:
                    error_text = response.text
                    if attempt < max_retries - 1:
                        wait_time = 3 * (2 ** attempt)  # 3s, 6s, 12s, 24s...
                        logger.warning(f"API 请求失败 ({response.status_code})，重试 {attempt + 1}/{max_retries} (等待 {wait_time}s)...")
                        time.sleep(wait_time)
                        continue
                    raise Exception(f"API 请求失败 ({response.status_code}): {error_text}")
                    
            except (requests.exceptions.SSLError, requests.exceptions.ConnectionError, ValueError) as e:
                # Catch ValueError (empty result) here to trigger retry
                if attempt < max_retries - 1:
                    wait_time = 3 * (2 ** attempt)
                    logger.warning(f"请求异常 ({type(e).__name__})，重试 {attempt + 1}/{max_retries} (等待 {wait_time}s)...: {e}")
                    time.sleep(wait_time)
                    continue
                raise
                
        # 如果所有重试都失败
        raise Exception("图片生成失败：所有重试均失败")
    
    def upscale_image(self, image_path: str, resolution: str = "4K") -> bool:
        """
        使用 Gemini API 高保真放大已有图片。
        只放大不改变任何排版、文字、颜色或设计元素。
        返回是否成功。
        """
        import requests
        from PIL import Image
        import io
        import base64
        import time

        resolution = resolution.upper()
        if resolution not in ("2K", "4K"):
            logger.warning(f"⚠️ 分辨率参数错误 ({resolution})，不支持放大，保持原图。")
            return False

        if not os.path.exists(image_path):
            logger.error(f"❌ 找不到图片文件: {image_path}")
            return False

        logger.info(f"正在高保真放大图片至 {resolution}: {image_path}")

        prompt = (
            f"Upscale this image to {resolution} resolution. ACT AS A HIGH-FIDELITY UPSCALER. "
            "You must maintain all text, details, layouts, and colors exactly as they appear in the source image. "
            "Do NOT change any words, do NOT move any text, do NOT add or remove any design elements. "
            "Simply increase the resolution, sharpness, and clarity."
        )

        with open(image_path, "rb") as f:
            image_bytes = f.read()
        
        # 兼容现有的请求组装逻辑
        mime_type = "image/png"
        if str(image_path).lower().endswith(('.jpg', '.jpeg')):
            mime_type = "image/jpeg"
            
        b64_data = base64.b64encode(image_bytes).decode("utf-8")
        parts = [
            {"text": prompt},
            {"inlineData": {"mimeType": mime_type, "data": b64_data}}
        ]

        # generationConfig
        generation_config = {
            "responseModalities": ["IMAGE"],
            "imageConfig": {
                "aspectRatio": "16:9",
                "imageSize": resolution
            }
        }
        
        payload = {
            "contents": [{"parts": parts}],
            "generationConfig": generation_config,
        }

        api_key = self.client.api_key if hasattr(self.client, 'api_key') else getattr(self, '_api_key', None)
        if not api_key:
            logger.error("❌ 无法获取 API Key")
            return False

        api_base = self.client.base_url if hasattr(self.client, 'base_url') else "https://generativelanguage.googleapis.com/v1beta/openai"
        # 针对 Gemini API，直接构造 REST URL (非 OpenAI 兼容 URL)
        # 如果提供了 openai base url，提取主机名并重组
        if "googleapis.com" in str(api_base):
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.image_model}:generateContent?key={api_key}"
        else:
            # 对于第三方反代，假设它是直接反代的
            base = str(api_base).replace("/openai/v1", "").replace("/openai", "")
            url = f"{base}/models/{self.image_model}:generateContent?key={api_key}"

        headers = {"Content-Type": "application/json"}
        
        # 重试逻辑
        max_retries = 5
        base_wait = 3
        for attempt in range(max_retries):
            try:
                response = requests.post(url, headers=headers, json=payload, timeout=180)
                if response.status_code == 200:
                    data = response.json()
                    candidates = data.get("candidates", [])
                    if candidates:
                        parts = candidates[0].get("content", {}).get("parts", [])
                        for part in parts:
                            inline_data = part.get("inlineData") or part.get("inline_data")
                            if inline_data:
                                img_bytes = base64.b64decode(inline_data["data"])
                                img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                                # 覆盖保存原图
                                img.save(image_path)
                                logger.info(f"✅ 成功放大图片并覆盖保存: {image_path}")
                                return True
                    logger.error(f"❌ API返回异常数据格式: {str(data)[:200]}...")
                elif response.status_code == 429:
                    wait_time = base_wait * (2 ** attempt)
                    logger.warning(f"⚠️ API 速率限制 (429)，等待 {wait_time} 秒后重试...")
                    time.sleep(wait_time)
                    continue
                else:
                    logger.error(f"❌ API请求失败 ({response.status_code}): {response.text}")
                    if attempt < max_retries - 1:
                        wait_time = base_wait * (2 ** attempt)
                        time.sleep(wait_time)
                        continue
                    break
            except Exception as e:
                logger.error(f"❌ 图片放大生成出错: {e}")
                if attempt < max_retries - 1:
                    wait_time = base_wait * (2 ** attempt)
                    time.sleep(wait_time)
                
        return False

    def create_pptx(self, images: List[Image.Image], output_path: str) -> str:
        """
        Legacy: 将图片列表转换为 PPTX 文件 (全屏直出模式)
        """
        logger.info(f"正在创建PPTX文件 (Legacy): {output_path}")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9 比例
        
        for i, image in enumerate(images):
            # 保存临时图片
            temp_path = self.output_dir / f"temp_slide_{i}.png"
            image.save(temp_path, "PNG")
            
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加图片 (全屏铺满)
            slide.shapes.add_picture(
                str(temp_path),
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            # 清理临时文件
            temp_path.unlink()
        
        prs.save(output_path)
        logger.info(f"PPTX文件已保存: {output_path}")
        return output_path

    def _calculate_dynamic_layout(self, bg_img: Image.Image, native_images: List[Dict]) -> List[Dict]:
        """Use Vision LLM to calculate perfect bounding boxes based on the actual generated background."""
        import copy
        import base64
        import json
        from io import BytesIO
        from openai import OpenAI
        
        # Deep copy to avoid mutating original plan dict globally if not needed
        updated_images = copy.deepcopy(native_images)
        
        try:
            buffered = BytesIO()
            bg_img.copy().convert("RGB").save(buffered, format="JPEG", quality=85)
            img_b64 = base64.b64encode(buffered.getvalue()).decode()
            
            img_details = []
            for idx, img_conf in enumerate(updated_images):
                path = img_conf.get('path')
                try:
                    from PIL import Image as PILImage
                    with PILImage.open(path) as n_img:
                        w, h = n_img.size
                        aspect = w / h
                    img_details.append(f"Image {idx+1}: Aspect Ratio {aspect:.2f} (Width/Height)")
                except Exception as e:
                    img_details.append(f"Image {idx+1}: Unknown aspect ratio (assume ~1.5)")

            prompt = f"""You are an expert presentation layout designer.
Look at the provided presentation slide image. This image already contains text and background graphics generated by another AI.
I need to place {len(updated_images)} native photo(s) on top of this slide. 

Your task is to find the ABSOLUTE PERFECT empty space (Safe Zone) on the slide to place the photo(s).
CRITICAL RULES:
1. Find the LARGEST possible rectangular empty space that DOES NOT OVERLAP ANY TEXT or important UI elements.
2. IMPORTANT: You must leave a generous padding/margin around your box. Do not let your box touch the text or the edges of the slide. If there is text on the left, start your box significantly to the right of it.
3. The bounding box must look visually balanced, aligning naturally with the text blocks (e.g., matching top/bottom margins).
4. Provide the coordinates as floats between 0.0 and 1.0 (where 0,0 is top-left, 1,1 is bottom-right).

Return ONLY a valid JSON array of bounding boxes in exactly the same order as the images listed above.
Example:
[
  {{"left": 0.55, "top": 0.20, "width": 0.40, "height": 0.60}}
]
"""
            api_key = self.client.api_key if hasattr(self.client, 'api_key') else getattr(self, '_api_key', None)
            
            client = OpenAI(api_key=api_key, base_url=self.client.base_url)
            
            logger.info(f"👁️ 正在使用视觉模型分析底层图片，寻找完美的排版位置...")
            response = client.chat.completions.create(
                model="gemini-3.1-pro-preview",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{img_b64}"
                                }
                            }
                        ]
                    }
                ],
                temperature=0.1
            )
            
            content = response.choices[0].message.content.strip()
            start = content.find('[')
            end = content.rfind(']')
            if start != -1 and end != -1:
                json_str = content[start:end+1]
                boxes = json.loads(json_str)
                
                for i, box in enumerate(boxes):
                    if i < len(updated_images):
                        updated_images[i]['dynamic_bounding_box'] = box
                        logger.info(f"  🎯 视觉模型计算出图 {i+1} 完美坐标: {box}")
                        
        except Exception as e:
            logger.warning(f"视觉排版分析失败，将回退到原始排版: {e}")
            
        return updated_images

    def create_advanced_pptx(self, visual_plan: List[Dict], images: Dict[int, Image.Image], output_path: str, template_path: str = None, project_dir: str = None) -> str:
        """
        Advanced: 将视觉计划和图片组装为 PPTX
        - 支持 .pptx 模版母版映射
        - 支持原生文本填充 (Editable Text)
        - 支持 Hex 颜色强制
        """
        logger.info(f"正在创建高级 PPTX 文件: {output_path}")
        
        if template_path and template_path.lower().endswith('.pptx') and os.path.exists(template_path):
            logger.info(f"📦 使用 .pptx 模版: {template_path}")
            prs = Presentation(template_path)
        else:
            logger.info("📄 使用默认空白模版")
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
        # Helper: Find layout by name
        def get_layout(prs, layout_name_hints):
            for layout in prs.slide_layouts:
                name = layout.name.lower()
                for hint in layout_name_hints:
                    if hint.lower() in name:
                        return layout
            return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0] # Fallback
            
        # Helper: Hex to RGB
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

        for i, slide_plan in enumerate(visual_plan):
            page_type = slide_plan.get('type', 'content').lower()
            text_content = slide_plan.get('text_content', {})
            style_config = slide_plan.get('style_config', {})
            table_data = slide_plan.get('table_data') or text_content.get('table_data')
            visualization = slide_plan.get('visualization', '')

            # 1. 对有全屏背景图的页面，一律使用 Blank 布局，避免未填充占位符显示为黑色块
            # 背景图/图表页：强制 Blank 布局，消除黑色占位符
            layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

            slide = prs.slides.add_slide(layout)

            page_num = slide_plan.get('page_num')
            img = images.get(page_num)

            if img:
                # 图表/普通页：添加背景图
                temp_path = self.output_dir / f"temp_slide_{page_num:02d}.png"
                img.save(temp_path, "PNG")
                slide.shapes.add_picture(str(temp_path), 0, 0, prs.slide_width, prs.slide_height)

            # Add Logo as a separate, movable PPTX shape (not burned into image)
            logo_path = slide_plan.get('logo_path')
            import os
            if logo_path and os.path.exists(logo_path) and page_type != 'background_only':
                logo_loc = (slide_plan.get('logo_location') or 'Top-Right').lower()
                logo_h = Inches(0.45)
                margin_x = Inches(0.3)
                margin_y = Inches(0.2)
                from PIL import Image as PILImage
                try:
                    _logo = PILImage.open(logo_path)
                    logo_aspect = _logo.width / _logo.height
                    logo_w = Inches(0.45 * logo_aspect)
                except Exception:
                    logo_w = Inches(1.2)

                if 'left' in logo_loc and 'top' in logo_loc:
                    lx, ly = margin_x, margin_y
                elif 'right' in logo_loc and 'bottom' in logo_loc:
                    lx = prs.slide_width - logo_w - margin_x
                    ly = prs.slide_height - logo_h - margin_y
                elif 'left' in logo_loc and 'bottom' in logo_loc:
                    lx, ly = margin_x, prs.slide_height - logo_h - margin_y
                else:
                    lx = prs.slide_width - logo_w - margin_x
                    ly = margin_y

                slide.shapes.add_picture(logo_path, lx, ly, logo_w, logo_h)

            # --- NEW: Add Multiple Native Images ---
            native_images = slide_plan.get('native_images', [])
            if not native_images and slide_plan.get('native_image'):
                native_images = [slide_plan.get('native_image')]

            # If we have native images and a generated background, use VLM to find exact placement
            if native_images and img:
                native_images = self._calculate_dynamic_layout(img, native_images)

            for img_conf in native_images:
                img_path = img_conf.get('path')
                
                # Check if it's an http path that hasn't been resolved yet
                if img_path and img_path.startswith("http"):
                    import os
                    from urllib.parse import urlparse
                    
                    # Try to find it relative to the original content file
                    content_file = slide_plan.get("style_config", {}).get("_meta", {}).get("content_file", "")
                    if not content_file:
                        # Sometimes passed globally via visual_plan root metadata
                        pass
                        
                    # Let's see if we can find it in the same directory as the script execution or project dir
                    filename = os.path.basename(urlparse(img_path).path)
                    
                    # Search up to find the document directory
                    possible_paths = [
                        filename,
                        os.path.join(project_dir, filename) if project_dir else filename,
                    ]
                    
                    # Add original document dir if available from the plan
                    # We can't access meta directly here easily, but we can search common locations
                    import glob
                    found = False
                    for search_path in [f"**/{filename}"]:
                        matches = glob.glob(search_path, recursive=True)
                        if matches:
                            img_path = matches[0]
                            found = True
                            break
                            
                    if not found:
                        logger.warning(f"无法找到网络图片对应的本地文件: {img_path}")
                        continue
                
                if not img_path or not os.path.exists(img_path):
                    logger.warning(f"原生图片不存在或路径为空: {img_path}")
                    continue
                    
                layout = img_conf.get('layout', 'center')
                bbox = img_conf.get('bounding_box')
                
                try:
                    from PIL import Image as PILImage
                    import tempfile
                    
                    # Convert WebP to PNG temporarily if needed since python-pptx doesn't support WebP
                    is_temp_file = False
                    if img_path.lower().endswith('.webp'):
                        try:
                            webp_img = PILImage.open(img_path)
                            temp_fd, temp_path = tempfile.mkstemp(suffix='.png')
                            os.close(temp_fd)
                            webp_img.save(temp_path, format="PNG")
                            img_path = temp_path
                            is_temp_file = True
                        except Exception as e:
                            logger.warning(f"无法转换 WebP 图片 {img_path}: {e}")
                            continue
                            
                    native_img = PILImage.open(img_path)
                    img_w, img_h = native_img.size
                    aspect = img_w / img_h
                    
                    sw = prs.slide_width
                    sh = prs.slide_height
                    
                    margin = Inches(0.5)
                    
                    # 1. Resolve target bounding box (left, top, max_width, max_height)
                    dynamic_bbox = img_conf.get('dynamic_bounding_box')
                    active_bbox = dynamic_bbox if dynamic_bbox else bbox
                    
                    if active_bbox:
                        # Semantic coordinate system (0.0 - 1.0)
                        target_l = sw * active_bbox.get('left', 0)
                        target_t = sh * active_bbox.get('top', 0)
                        max_w = sw * active_bbox.get('width', 1.0)
                        max_h = sh * active_bbox.get('height', 1.0)
                        logger.info(f"    计算原生图片坐标: target_l={target_l}, target_t={target_t}, max_w={max_w}, max_h={max_h} (来自 VLM: {bool(dynamic_bbox)})")
                    else:
                        # Legacy enum system
                        if layout == 'right_half':
                            box = (sw / 2 + margin/2, margin, sw / 2 - margin*1.5, sh - margin*2)
                        elif layout == 'left_half':
                            box = (margin, margin, sw / 2 - margin*1.5, sh - margin*2)
                        elif layout == 'bottom_right':
                            box = (sw * 0.6, sh * 0.5, sw * 0.4 - margin, sh * 0.5 - margin)
                        elif layout == 'fullscreen':
                            box = (0, 0, sw, sh)
                        else: # center
                            box = (margin*2, margin*2, sw - margin*4, sh - margin*4)
                        target_l, target_t, max_w, max_h = box
                        
                    # 2. Calculate fitted dimensions preserving aspect ratio
                    if max_h == 0: max_h = 1 # prevent div by zero
                    target_aspect = max_w / max_h
                    if aspect > target_aspect:
                        # Image is wider than target box, so width is the limiting factor
                        final_w = max_w
                        final_h = max_w / aspect
                    else:
                        # Image is taller than target box, so height is the limiting factor
                        final_h = max_h
                        final_w = max_h * aspect
                        
                    # 3. Align within the target box
                    if dynamic_bbox:
                        # VLM calculates the safe zone. We MUST perfectly center the image inside this safe zone.
                        # This prevents the image from sticking to the edge and maintains the VLM's intended margins.
                        final_l = target_l + (max_w - final_w) / 2
                        final_t = target_t + (max_h - final_h) / 2
                    else:
                        # Old mechanical logic fallback
                        if active_bbox:
                            left_pct = active_bbox.get('left', 0)
                            if left_pct < 0.2:
                                # Align left
                                final_l = target_l
                            elif left_pct > 0.5:
                                # Align right
                                final_l = target_l + (max_w - final_w)
                            else:
                                # Center horizontal
                                final_l = target_l + (max_w - final_w) / 2
                        else:
                            final_l = target_l + (max_w - final_w) / 2
                            
                        # Always center vertically for now
                        final_t = target_t + (max_h - final_h) / 2
                        
                    slide.shapes.add_picture(img_path, final_l, final_t, final_w, final_h)
                    logger.info(f"  已插入多图排版图片: {active_bbox if active_bbox else layout}")
                    
                    # Clean up temp file if we created one
                    if is_temp_file:
                        try:
                            os.remove(img_path)
                        except:
                            pass
                            
                except Exception as e:
                    logger.warning(f"无法插入多图排版图片 {img_path}: {e}")
                    if 'is_temp_file' in locals() and is_temp_file:
                        try:
                            os.remove(img_path)
                        except:
                            pass

            # 添加演讲者备注 (Speaker Notes)
            speaker_notes = slide_plan.get('speaker_notes')
            if speaker_notes:
                try:
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = speaker_notes
                except Exception as e:
                    logger.warning(f"无法添加演讲者备注到第 {i+1} 页: {e}")

        prs.save(output_path)
        logger.info(f"高级 PPTX 文件已保存: {output_path}")
        return output_path

if __name__ == "__main__":
    pass
