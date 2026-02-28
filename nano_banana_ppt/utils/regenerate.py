"""
重新生成失败的PPT页面
"""
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

# 添加路径以便导入
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '../../..')))
from nano_banana_ppt.core.generator import PPTGenerator

def regenerate_failed_pages(pptx_file: str, plan_file: str, failed_pages: list = None):
    """
    重新生成失败的页面并更新PPT
    需配合 visual_plan.json 使用
    """
    import json
    
    if not os.path.exists(plan_file):
        print("❌ 找不到 visual_plan.json")
        return

    with open(plan_file, 'r', encoding='utf-8') as f:
        visual_plan = json.load(f)
    
    # 建立页码映射
    pages_dict = {p['page_num']: p for p in visual_plan}
    
    # 打开现有PPT
    prs = Presentation(pptx_file)
    
    # 自动检测失败页面
    if failed_pages is None:
        failed_pages = []
        for i, slide in enumerate(prs.slides, 1):
            # 检查是否是占位图
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    try:
                        img_bytes = shape.image.blob
                        img = Image.open(io.BytesIO(img_bytes))
                        colors = img.getcolors(maxcolors=256*256*256)
                        if colors and len(colors) < 10:  # 颜色少视为占位图
                            failed_pages.append(i)
                    except:
                        pass
                    break
    
    if not failed_pages:
        print("✅ 没有发现失败的页面")
        return
    
    print(f"发现 {len(failed_pages)} 页需要重新生成: {failed_pages}")
    
    api_key = os.getenv("OPENAI_API_KEY")
    api_base = os.getenv("OPENAI_API_BASE")
    generator = PPTGenerator(api_key, api_base)
    
    # 重新生成
    for page_num in failed_pages:
        if page_num not in pages_dict:
            continue
            
        page = pages_dict[page_num]
        print(f"\n正在重新生成第 {page_num} 页...")
        
        try:
            image = generator.generate_image(page['visual_prompt'], aspect_ratio="16:9")
            
            temp_path = generator.output_dir / f"temp_regen_{page_num}.png"
            image.save(temp_path, "PNG")
            
            # 更新PPT (简单替换逻辑：删除旧slide，插入新slide会导致顺序问题，这里仅示意)
            # 更好的做法是更新 image blob，或者重新组装整个 PPT
            # 暂时只打印提示，完整实现需操作 pptx 底层
            print(f"✅ 图片已生成: {temp_path}")
            print("⚠️ 请手动替换 PPT 中的图片，或重新运行生成流程。")
            
        except Exception as e:
            print(f"❌ 重新生成失败: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法: python -m nano_banana_ppt.utils.regenerate <pptx_file> <plan_json> [pages]")
        sys.exit(1)
    
    pptx = sys.argv[1]
    plan = sys.argv[2]
    pages = [int(p) for p in sys.argv[3:]] if len(sys.argv) > 3 else None
    
    regenerate_failed_pages(pptx, plan, pages)
