"""
PM Agent - 产品经理 Agent

负责接住用户输入并路由到正确阶段，是系统的入口和协调者。
"""

import os
import logging
import json
from typing import Optional, Dict, Any, List
from pathlib import Path

from ..utils.image_assets import ImageAssetsManager, ImageAsset, ImageMode
from ..utils.provider_config import DEFAULT_LLM_MODEL

logger = logging.getLogger(__name__)

# 固定的输出目录
DEFAULT_OUTPUT_BASE = Path.home() / "Desktop" / "AI output" / "ppt"


class PMAgent:
    """PM Agent - 负责用户输入接收、意图理解和流程路由"""

    def __init__(self, client, project_dir: str = "output"):
        """
        初始化 PM Agent

        Args:
            client: OpenAI 兼容客户端
            project_dir: 项目输出目录
        """
        self.client = client
        self.project_dir = Path(project_dir)
        self.project_dir.mkdir(parents=True, exist_ok=True)

        # 初始化管理器
        self.image_assets_manager = ImageAssetsManager(
            str(self.project_dir / "image_assets.json"),
            client=client
        )

    def intake(self, user_input: Dict[str, Any]) -> Dict[str, Any]:
        """
        接收用户输入，进行初步处理

        Args:
            user_input: 用户输入，可能包含：
                - text: 文本内容
                - images: 图片路径列表
                - template_pptx: 毛坯 PPT 路径
                - urls: 网络图片 URL 列表

        Returns:
            处理结果，包含下一步建议
        """
        logger.info("📥 PM Agent 接收用户输入")

        result = {
            "status": "success",
            "next_gate": None,
            "message": "",
        }

        try:
            # 1. 提取图片资产
            if user_input.get("images") or user_input.get("urls") or user_input.get("template_pptx"):
                self._extract_images_from_input(user_input)

            # 2. 整理用户意图
            if user_input.get("text"):
                self._organize_user_intent(user_input["text"])

            # 3. 判断当前应进入哪个 Gate
            next_gate = self.determine_gate()
            result["next_gate"] = next_gate
            result["message"] = f"建议进入 {next_gate} 阶段"

            return result

        except Exception as e:
            logger.error(f"PM Agent 处理失败: {e}")
            result["status"] = "error"
            result["message"] = str(e)
            return result

    def update_state(self, updates: Dict[str, Any]) -> Dict[str, Any]:
        """
        更新项目状态（供 Claude Code 调用）

        Args:
            updates: 状态更新字典，支持的 key：
                - logo_path: Logo 文件路径
                - style_guide_file: 设计指导文件路径
                - style_preference: 风格偏好（文本）
                - visual_constraints: 视觉约束列表（如 ["插图使用写实风格"]）
                - briefing: 用户意图
                - user_confirmed_visual_plan: bool，用户是否已确认 visual_plan

        Returns:
            {"status": "success", "message": "状态已更新"}
        """
        state = ProjectState(str(self.project_dir))

        self._update_assets(state, updates)
        self._update_constraints(state, updates)
        self._update_confirmation(state, updates)

        state.persist()
        return {"status": "success", "message": "状态已更新"}

    def _update_assets(self, state, updates: Dict[str, Any]) -> None:
        """更新 assets（logo 等）"""
        if "logo_path" not in updates:
            return

        logo_path = updates["logo_path"]
        if os.path.exists(logo_path):
            assets = state.get("assets", {})
            assets["logo_path"] = logo_path
            state.set("assets", assets)

            # 同步更新 meta
            meta = state.get("meta", {})
            meta["logo_file"] = logo_path
            state.set("meta", meta)
            logger.info(f"✅ Logo 已更新: {logo_path}")
        else:
            logger.warning(f"⚠️ Logo 文件不存在: {logo_path}")

    def _update_constraints(self, state, updates: Dict[str, Any]) -> None:
        """更新 constraints（风格、约束等）"""
        constraints = state.get("constraints", {})

        if "style_preference" in updates:
            constraints["style_preference"] = updates["style_preference"]

        if "style_guide_file" in updates:
            guide_file = updates["style_guide_file"]
            if os.path.exists(guide_file):
                with open(guide_file, 'r', encoding='utf-8') as f:
                    guide_content = f.read()
                constraints["style_guide_content"] = guide_content
                constraints["style_guide_file"] = guide_file
                logger.info(f"✅ 设计指导已加载: {guide_file}")
            else:
                logger.warning(f"⚠️ 设计指导文件不存在: {guide_file}")

        if "visual_constraints" in updates:
            constraints["visual_constraints"] = updates["visual_constraints"]

        if "briefing" in updates:
            constraints["briefing"] = updates["briefing"]

        state.set("constraints", constraints)

    def _update_confirmation(self, state, updates: Dict[str, Any]) -> None:
        """更新确认状态"""
        if "user_confirmed_visual_plan" in updates:
            state.set("user_confirmed_visual_plan", updates["user_confirmed_visual_plan"])

    def execute_phase(self, phase: str, **kwargs) -> Dict[str, Any]:
        """
        执行指定阶段的完整流程。

        Args:
            phase: 阶段名称 ("content" | "visual" | "execute")
            **kwargs: 阶段特定参数

        Returns:
            阶段执行结果
        """
        logger.info(f"🚀 PM Agent 开始执行阶段: {phase}")

        if phase == "content":
            return self._execute_content_phase(**kwargs)
        elif phase == "visual":
            return self._execute_visual_phase(**kwargs)
        elif phase == "execute":
            return self._execute_execute_phase(**kwargs)
        else:
            return {"status": "error", "message": f"未知阶段: {phase}"}

    def _execute_content_phase(
        self,
        content_file: str = None,
        template_file: str = None,
        logo_file: str = None,
        page_count: int = None,
        briefing: str = None,
        output_name: str = None,
    ) -> Dict[str, Any]:
        """
        执行内容阶段：调用 NarrativeAgent 生成内容大纲。

        Args:
            content_file: 内容文件路径
            template_file: 模板文件路径
            logo_file: Logo 文件路径
            page_count: 期望页数
            briefing: 用户意图
            output_name: 输出名称

        Returns:
            执行结果
        """
        from datetime import date
        import re

        from ..agents.narrative import NarrativeAgent
        from ..agents.template import TemplateAgent
        from ..utils.review_plan import build_content_review_md
        from ..utils.compile_content_plan import compile_content_plan
        from ..utils.provider_config import get_llm_api_key, get_llm_api_base

        result = {"status": "success", "message": ""}

        # 1. 确定项目目录（内联 _resolve_project_dir 逻辑以避免循环导入）
        if content_file:
            date_prefix = date.today().strftime("%Y%m%d")
            if output_name:
                name = output_name
            else:
                content_path = Path(content_file)
                parent_parts = content_path.parts
                for i, part in enumerate(parent_parts):
                    if i > 0 and parent_parts[i-1] == "ppt" and part.startswith("202"):
                        name_part = "_".join(parent_parts[i:]).replace("_content_plan", "").replace("_content", "")
                        name_part = re.sub(r'^\d{8}_', '', name_part)
                        name_part = re.sub(r'^\d{8}_\d{8}_', '', name_part)
                        if name_part:
                            name = name_part
                            break
                else:
                    name = content_path.stem
                    name = re.sub(r'_content_plan$', '', name)
                    name = re.sub(r'_content$', '', name)

            if re.match(r'^\d{8}_', name):
                dir_name = name
            else:
                dir_name = f"{date_prefix}_{name}"

            project_dir = DEFAULT_OUTPUT_BASE / dir_name
            project_dir.mkdir(parents=True, exist_ok=True)
            self.project_dir = project_dir

        # 2. 初始化状态管理器
        state = ProjectState(str(self.project_dir))

        # 3. 读取内容
        if not content_file or not os.path.exists(content_file):
            return {"status": "error", "message": f"内容文件不存在: {content_file}"}

        if content_file.lower().endswith('.pdf'):
            try:
                import fitz
                doc = fitz.open(content_file)
                parts = []
                for page in doc:
                    parts.append(page.get_text())
                doc.close()
                content_context = "\n\n".join(parts)
            except Exception as e:
                return {"status": "error", "message": f"PDF 读取失败: {e}"}
        else:
            with open(content_file, 'r', encoding='utf-8') as f:
                content_context = f.read()

        # 4. 初始化 Agent
        api_key = get_llm_api_key()
        api_base = get_llm_api_base()
        narrative_agent = NarrativeAgent(api_key, api_base, project_dir=str(self.project_dir))
        template_agent = TemplateAgent(api_key, api_base, output_dir=str(self.project_dir / "template_assets"))

        # 5. 分析内容
        logger.info("🔍 [Content Phase] 分析内容...")
        inferred = narrative_agent.analyze_content(content_context)
        state.set("inferred", inferred)

        # 6. 解析模板
        template_info = None
        assets = {}
        if template_file and os.path.exists(template_file):
            logger.info("🖼️ [Content Phase] 解析模板...")
            try:
                template_info = template_agent.process_template(template_file)
                assets['template_file'] = template_file
                assets['logo_path'] = template_info.get('logo_path')
                assets['template_images'] = template_info.get('reference_images')
            except Exception as e:
                logger.warning(f"模板解析失败: {e}")

        # 7. Logo 处理
        if logo_file and os.path.exists(logo_file):
            assets['logo_path'] = logo_file
        elif logo_file:
            logger.warning(f"Logo 文件不存在: {logo_file}")

        # 8. 品牌色提取
        if assets.get('logo_path') and not template_info:
            from ..utils.image_utils import extract_dominant_colors
            brand_colors = extract_dominant_colors(assets['logo_path'], num_colors=3)
            if brand_colors:
                inferred['brand_colors'] = brand_colors

        # 9. 构建约束
        constraints = {
            "target_audience": inferred.get("target_audience", "通用受众"),
            "presentation_type": inferred.get("presentation_type", "商业演示"),
            "duration": inferred.get("duration", "15分钟"),
            "style_preference": inferred.get("style_preference", "专业严谨"),
            "page_count": str(page_count) if page_count else "10",
            "briefing": briefing,
        }
        state.set("constraints", constraints)
        state.set("assets", assets)

        # 10. 生成叙事大纲
        logger.info("📝 [Content Phase] 生成叙事大纲...")
        narrative_outline = narrative_agent.generate_narrative_outline(
            content_context, constraints, content_file_path=content_file
        )
        state.set("narrative_outline", narrative_outline)

        # 11. 保存元数据
        meta = {
            "project_name": self.project_dir.name,
            "project_dir": str(self.project_dir),
            "content_file": content_file,
            "template_file": assets.get('template_file'),
            "logo_file": assets.get('logo_path'),
        }
        state.set("meta", meta)

        # 12. 持久化状态
        state.persist()

        # 13. 生成 content_plan.md
        content_md_content = build_content_review_md(narrative_outline, meta)
        content_md_path = self.project_dir / "content_plan.md"
        with open(content_md_path, 'w', encoding='utf-8') as f:
            f.write(content_md_content)

        # 14. 编译 content_plan.json
        compile_content_plan(str(content_md_path), str(self.project_dir / "content_plan.json"))

        logger.info(f"✅ 内容阶段完成: {len(narrative_outline)} 页")
        result["message"] = f"内容大纲已生成，共 {len(narrative_outline)} 页"
        result["content_plan_path"] = str(content_md_path)
        result["page_count"] = len(narrative_outline)

        return result

    def _execute_visual_phase(
        self,
        project_dir: str = None,
        style_preference: str = None,
    ) -> Dict[str, Any]:
        """
        执行视觉阶段：调用 VisualAgent 生成视觉计划。

        Args:
            project_dir: 项目目录
            style_preference: 风格偏好

        Returns:
            执行结果
        """
        from ..agents.visual import VisualAgent
        from ..utils.review_plan import (
            build_visual_plan_from_content_plan,
            parse_review_md,
            derive_technical_plan,
            generate_per_slide_visual_suggestions,
            generate_design_manifesto,
            REVIEW_MD_FILENAME,
        )
        from ..utils.provider_config import get_llm_api_key, get_llm_api_base

        result = {"status": "success", "message": ""}

        # 1. 确定项目目录
        if project_dir:
            self.project_dir = Path(project_dir)

        # 2. 加载状态
        state = ProjectState(str(self.project_dir))
        narrative_outline = state.get("narrative_outline", [])
        assets = state.get("assets", {})
        meta = state.get("meta", {})
        template_info = state.get("template_info")
        inferred = state.get("inferred", {})
        constraints = state.get("constraints", {})

        # 3. 检查是否有用户编辑过的 content_plan.md
        content_md_path = self.project_dir / "content_plan.md"
        if content_md_path.exists():
            from ..utils.doc_normalizer import normalize_content_plan
            with open(content_md_path, 'r', encoding='utf-8') as f:
                md_text = f.read()

            normalized_content, issues = normalize_content_plan(md_text)
            if normalized_content != md_text:
                with open(content_md_path, 'w', encoding='utf-8') as f:
                    f.write(normalized_content)
                logger.info("✅ content_plan.md 已规范化")

            parsed = parse_review_md(normalized_content)
            if parsed and parsed.get("pages"):
                narrative_outline = parsed.get("pages")
                state.set("narrative_outline", narrative_outline)

        # 4. 初始化 Agent
        api_key = get_llm_api_key()
        api_base = get_llm_api_base()
        prompt_mode = os.getenv("PROMPT_MODE", "verbose")
        visual_agent = VisualAgent(api_key, api_base, prompt_mode=prompt_mode)

        # 5. 风格定义
        logger.info("🎨 [Visual Phase] 定义视觉风格...")

        # 优先使用用户通过 update_state() 提供的风格偏好
        if not constraints.get("style_preference"):
            constraints["style_preference"] = style_preference or inferred.get("style_preference", "专业商务")

        # 如果有设计指导文件，注入到 constraints
        if constraints.get("style_guide_content"):
            logger.info(f"📄 使用设计指导: {constraints.get('style_guide_file')}")
            # Visual Agent 会在 define_style() 中读取 style_guide_content

        brand_colors = inferred.get("brand_colors", [])
        if brand_colors:
            constraints["brand_colors"] = brand_colors

        style_definition = visual_agent.define_style(constraints, assets, template_info)
        if isinstance(style_definition, tuple):
            style_desc_str, style_config = style_definition
        else:
            style_desc_str = str(style_definition)
            style_config = {"description": style_desc_str, "palette": [], "mode": "ai_minting"}

        state.set("style_config", style_config)

        # 6. 生成每页视觉描述
        logger.info("📋 [Visual Phase] 生成每页视觉描述...")

        # 传递用户的视觉约束（如"插图使用写实风格"）
        visual_constraints = constraints.get("visual_constraints", [])

        per_slide_descriptions = generate_per_slide_visual_suggestions(
            narrative_outline, style_config, api_key, api_base,
            visual_constraints=visual_constraints  # 新增参数
        )

        for page in narrative_outline:
            page_num = page.get("page_num")
            if page_num in per_slide_descriptions:
                page["visual_description"] = per_slide_descriptions[page_num]
                page["visual_suggestion"] = per_slide_descriptions[page_num]

        # 7. 生成 Design Manifesto
        logger.info("📋 [Visual Phase] 生成 Design Manifesto...")
        is_template_mode = bool(meta.get("template_file"))
        parsed_stub = {"pages": narrative_outline, "style": style_config}
        manifesto_dict = generate_design_manifesto(
            parsed=parsed_stub,
            template_mode=is_template_mode,
            client=visual_agent.client
        )
        manifesto_text = manifesto_dict.get("chinese_proposal", "")

        state.set("manifesto_bans", manifesto_dict.get("english_cliche_bans", ""))
        state.set("visual_diversity_strategy", manifesto_dict.get("visual_diversity_strategy", ""))

        # 8. 生成最终视觉计划
        logger.info("🧠 [Visual Phase] 生成最终视觉计划...")
        style_config_with_manifesto = dict(style_config)
        style_config_with_manifesto["manifesto"] = manifesto_text

        generated_slides = visual_agent.generate_visual_plan(
            narrative_outline=narrative_outline,
            style_definition_tuple=(style_config_with_manifesto.get("description", style_desc_str), style_config_with_manifesto),
            assets={"logo_path": meta.get("logo_file")},
            template_info=template_info,
            meta=meta,
        )

        # 9. 生成 visual_plan.md
        logger.info("📋 [Visual Phase] 生成 visual_plan.md...")
        review_md_path = self.project_dir / REVIEW_MD_FILENAME

        review_md_content = build_visual_plan_from_content_plan(
            str(content_md_path),
            style_config,
            meta,
            manifesto=manifesto_text,
            per_slide_descriptions=per_slide_descriptions,
            state_narrative_outline=narrative_outline,
            generated_slides=generated_slides,
        )
        with open(review_md_path, 'w', encoding='utf-8') as f:
            f.write(review_md_content)

        # 10. 生成 visual_plan.json
        parsed_review = parse_review_md(review_md_content, project_dir=str(self.project_dir))
        visual_plan_data = derive_technical_plan(
            parsed_review,
            str(self.project_dir),
            meta.get("content_file", str(content_md_path)),
            api_key,
            api_base,
        )
        visual_plan_json_path = self.project_dir / "visual_plan.json"
        with open(visual_plan_json_path, 'w', encoding='utf-8') as f:
            json.dump(visual_plan_data, f, ensure_ascii=False, indent=2)

        # 11. 持久化状态
        state.persist()

        logger.info(f"✅ 视觉阶段完成: {len(generated_slides)} 页")
        result["message"] = f"视觉计划已生成，共 {len(generated_slides)} 页"
        result["visual_plan_path"] = str(review_md_path)

        return result

    def _execute_execute_phase(
        self,
        plan_file: str = None,
        output_name: str = None,
        resolution: str = "1K",
        slide_filter: list = None,
        reassemble_only: bool = False,
        force: bool = False,
    ) -> Dict[str, Any]:
        """
        执行阶段：调用 executor 生成最终 PPT。

        Args:
            plan_file: 计划文件路径或项目目录
            output_name: 输出名称
            resolution: 分辨率
            slide_filter: 仅执行的页面
            reassemble_only: 仅重新组装
            force: 强制执行

        Returns:
            执行结果
        """
        from ..core.executor import execute_plan
        from ..utils.review_plan import REVIEW_MD_FILENAME
        from ..utils.provider_config import get_llm_api_key, get_llm_api_base

        result = {"status": "success", "message": ""}

        # 1. 解析输入
        if plan_file:
            plan_path = Path(plan_file)
            if plan_path.is_dir():
                # 优先使用 visual_plan.json（已确认的技术计划）
                if (plan_path / "visual_plan.json").exists():
                    plan_file = str(plan_path / "visual_plan.json")
                    project_dir = str(plan_path)
                elif (plan_path / REVIEW_MD_FILENAME).exists():
                    plan_file = str(plan_path / REVIEW_MD_FILENAME)
                    project_dir = str(plan_path)
                else:
                    return {"status": "error", "message": f"目录中未找到计划文件: {plan_path}"}
            else:
                # 如果传入的是文件路径，project_dir 是文件所在目录
                project_dir = str(plan_path.parent)
            self.project_dir = Path(project_dir)

        # 2. 执行最终意图审查（智能跳过）
        state = ProjectState(str(self.project_dir))
        user_confirmed = state.get("user_confirmed_visual_plan", False)

        if not force and not user_confirmed:
            logger.info("🔍 执行最终意图审查...")
            review_report = self.final_intent_review()
            if not review_report["passed"]:
                result["status"] = "error"
                result["message"] = f"最终审查未通过: {review_report['summary']['blocking_issues']} 个阻塞问题"
                return result
        else:
            if user_confirmed:
                logger.info("⏭️ 用户已确认 visual_plan，跳过审查")
            else:
                logger.info("⏭️ 使用 --force 参数，跳过审查")

        # 3. 执行生成
        logger.info("⚡ [Execute Phase] 开始生成 PPT...")
        logger.info(f"📄 读取计划文件: {plan_file}")

        with open(plan_file, 'r', encoding='utf-8') as f:
            plan_data = json.load(f)

        slides = plan_data if isinstance(plan_data, list) else plan_data.get("slides", plan_data)
        name = output_name or plan_data.get("meta", {}).get("project_name", self.project_dir.name)

        exec_plan_file = str(self.project_dir / "_exec_slides.json")
        with open(exec_plan_file, 'w', encoding='utf-8') as f:
            json.dump(slides, f, ensure_ascii=False, indent=2)

        out_path, _ = execute_plan(
            exec_plan_file,
            name,
            template_path=plan_data.get("meta", {}).get("template_file"),
            project_dir=str(self.project_dir),
            resolution=resolution,
            slide_filter=slide_filter,
            reassemble_only=reassemble_only,
        )

        if os.path.exists(exec_plan_file):
            os.remove(exec_plan_file)

        result["message"] = "PPT 生成完成"
        result["output_path"] = str(out_path)

        return result

    def determine_gate(self) -> str:
        """
        判断当前应进入哪个 Gate

        Returns:
            Gate 名称：Content / Visual / Execute
        """
        # 检查是否有 content_plan
        content_plan_path = self.project_dir / "content_plan.md"
        if not content_plan_path.exists():
            return "Content"

        # 检查是否有 visual_plan
        visual_plan_path = self.project_dir / "visual_plan.md"
        if not visual_plan_path.exists():
            return "Visual"

        # 都有了，进入执行阶段
        return "Execute"

    def _extract_images_from_input(self, user_input: Dict[str, Any]) -> None:
        """
        从用户输入中抽取图片资产

        Args:
            user_input: 用户输入
        """
        logger.info("🖼️ 开始抽取图片资产")

        # 1. 从本地路径列表抽取图片
        if user_input.get("images"):
            self._extract_from_local_paths(user_input["images"])

        # 2. 从 URL 列表下载图片
        if user_input.get("urls"):
            self._extract_from_urls(user_input["urls"])

        # 3. 从毛坯 PPTX 中提取图片
        if user_input.get("template_pptx"):
            self._extract_from_pptx(user_input["template_pptx"])

        # 4. 保存图片资产
        self.image_assets_manager.save_to_json()
        logger.info(f"✅ 图片资产抽取完成，共 {len(self.image_assets_manager.assets)} 张图片")

    def _extract_from_local_paths(self, image_paths: List[str]) -> None:
        """
        从本地路径列表抽取图片

        Args:
            image_paths: 图片路径列表
        """
        for path in image_paths:
            path_obj = Path(path)
            if not path_obj.exists():
                logger.warning(f"图片不存在: {path}")
                continue

            # 检查是否已存在
            if self.image_assets_manager.get_asset_by_path(path):
                logger.info(f"图片已存在，跳过: {path}")
                continue

            # 创建图片资产
            asset = ImageAsset(path=path)
            self.image_assets_manager.add_asset(asset)

            # 执行轻量读图
            try:
                result = self.image_assets_manager.light_read_image(path)
                if result:
                    logger.info(f"✅ 轻量读图完成: {path} -> {result.get('image_type')}")
                else:
                    logger.warning(f"轻量读图失败: {path}")
            except Exception as e:
                logger.error(f"轻量读图异常: {path} ({e})")

    def _extract_from_urls(self, urls: List[str]) -> None:
        """
        从 URL 列表下载图片

        Args:
            urls: 图片 URL 列表
        """
        import requests
        from urllib.parse import urlparse
        import hashlib

        for url in urls:
            try:
                # 下载图片
                response = requests.get(url, timeout=30)
                response.raise_for_status()

                # 生成文件名（使用 URL hash）
                url_hash = hashlib.md5(url.encode()).hexdigest()[:8]
                parsed = urlparse(url)
                ext = Path(parsed.path).suffix or ".jpg"
                filename = f"downloaded_{url_hash}{ext}"
                save_path = self.project_dir / "images" / filename

                # 确保目录存在
                save_path.parent.mkdir(parents=True, exist_ok=True)

                # 保存图片
                with open(save_path, 'wb') as f:
                    f.write(response.content)

                logger.info(f"✅ 下载图片: {url} -> {save_path}")

                # 创建图片资产
                asset = ImageAsset(path=str(save_path))
                self.image_assets_manager.add_asset(asset)

                # 执行轻量读图
                try:
                    result = self.image_assets_manager.light_read_image(str(save_path))
                    if result:
                        logger.info(f"✅ 轻量读图完成: {save_path} -> {result.get('image_type')}")
                except Exception as e:
                    logger.error(f"轻量读图异常: {save_path} ({e})")

            except Exception as e:
                logger.error(f"下载图片失败: {url} ({e})")

    def _extract_from_pptx(self, pptx_path: str) -> None:
        """
        从毛坯 PPTX 中提取图片

        Args:
            pptx_path: PPTX 文件路径
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("需要安装 python-pptx: pip install python-pptx")
            return

        try:
            prs = Presentation(pptx_path)
            image_count = 0

            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    # 检查是否是图片
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob

                        # 生成文件名
                        ext = image.ext
                        filename = f"pptx_slide{slide_idx + 1}_img{image_count + 1}.{ext}"
                        save_path = self.project_dir / "images" / filename

                        # 确保目录存在
                        save_path.parent.mkdir(parents=True, exist_ok=True)

                        # 保存图片
                        with open(save_path, 'wb') as f:
                            f.write(image_bytes)

                        logger.info(f"✅ 提取图片: 第 {slide_idx + 1} 页 -> {save_path}")

                        # 创建图片资产
                        asset = ImageAsset(path=str(save_path))
                        self.image_assets_manager.add_asset(asset)

                        # 执行轻量读图
                        try:
                            result = self.image_assets_manager.light_read_image(str(save_path))
                            if result:
                                logger.info(f"✅ 轻量读图完成: {save_path} -> {result.get('image_type')}")
                        except Exception as e:
                            logger.error(f"轻量读图异常: {save_path} ({e})")

                        image_count += 1

            logger.info(f"✅ 从 PPTX 提取了 {image_count} 张图片")

        except Exception as e:
            logger.error(f"从 PPTX 提取图片失败: {pptx_path} ({e})")

    def _organize_user_intent(self, text: str) -> None:
        """
        将用户原始输入整理为结构化的任务约束

        Args:
            text: 用户文本输入
        """
        logger.info("📝 开始整理用户意图")

        # 构建 prompt
        prompt = f"""你是一个专业的产品经理，需要将用户的原始输入整理为结构化的任务约束。

用户输入：
{text}

请分析用户输入，提取以下信息（以 JSON 格式返回）：

1. input_type: 输入类型，从以下选项中选择一个
   - "topic": 只给了一个主题或标题
   - "article": 提供了完整的文章或长文本
   - "outline": 提供了大纲或结构化内容
   - "template_ppt": 基于已有的毛坯 PPT 进行修改
   - "add_images": 为已有内容补充图片
   - "modify_slide": 修改单个或少数几页

2. goal: 任务目标（一句话概括用户想做什么）

3. audience: 目标受众（如果用户提到了，否则为 null）

4. tone: 语气/风格（如果用户提到了，否则为 null）
   例如：专业、轻松、学术、商务等

5. constraints: 关键限制（列表，提取用户明确提到的约束）
   例如：页数限制、时间限制、必须包含的内容等

6. image_anchors: 图片位置意图（列表，提取用户提到的图片需求）
   每个元素包含：
   - anchor: 语义锚点（图片应该出现在哪个内容段落）
   - description: 图片描述或要求

返回格式：
{{
  "input_type": "...",
  "goal": "...",
  "audience": "..." or null,
  "tone": "..." or null,
  "constraints": [...],
  "image_anchors": [
    {{"anchor": "...", "description": "..."}},
    ...
  ]
}}

只返回 JSON，不要有其他内容。"""

        try:
            # 调用 LLM
            response = self.client.chat.completions.create(
                model=DEFAULT_LLM_MODEL,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
            )

            result_text = response.choices[0].message.content.strip()

            # 提取 JSON（去除可能的 markdown 包裹）
            import json
            if result_text.startswith("```"):
                result_text = result_text.split("```")[1]
                if result_text.startswith("json"):
                    result_text = result_text[4:]
            result_text = result_text.strip()

            result = json.loads(result_text)

            logger.info(f"✅ 意图分析完成: {result['input_type']}")

            # Brief 功能已移除，意图分析结果仅用于日志记录

        except Exception as e:
            logger.error(f"整理用户意图失败: {e}")
            # 直接抛出异常，不再使用简化 Brief 降级
            raise RuntimeError(f"PM Agent 整理用户意图失败: {e}")

    def final_intent_review(self) -> Dict[str, Any]:
        """
        执行前的最终意图审查

        检查项：
        1. 资源完整性（调用 resource_check）
        2. Prompt 质量（调用 prompt_guard）
        3. Brief 与 visual_plan 的一致性

        Returns:
            审查报告，包含：
            - passed: 是否通过审查
            - resource_check: 资源检查报告
            - prompt_review: Prompt 审查报告
            - consistency_check: 一致性检查结果
            - summary: 汇总信息
        """
        from ..utils.resource_check import check_resources, print_report as print_resource_report
        from ..utils.prompt_guard import review_visual_plan, print_review_report
        import json

        logger.info("🔍 开始最终意图审查")

        report = {
            "passed": False,
            "resource_check": None,
            "prompt_review": None,
            "consistency_check": None,
            "summary": {
                "total_issues": 0,
                "blocking_issues": 0,
                "warnings": 0,
            }
        }

        # 1. 资源完整性检查
        logger.info("📋 检查资源完整性...")
        try:
            resource_report = check_resources(
                output_dir=str(self.project_dir),
                check_images=True,
                check_placeholders=True,
                check_fields=True,
                check_consistency=True
            )
            report["resource_check"] = resource_report.to_dict()

            # 统计问题
            report["summary"]["total_issues"] += len(resource_report.issues)
            report["summary"]["blocking_issues"] += resource_report.summary.get("error", 0)
            report["summary"]["warnings"] += resource_report.summary.get("warning", 0)

            if not resource_report.passed:
                logger.warning(f"⚠️ 资源检查未通过：{resource_report.summary.get('error', 0)} 个错误")
                print_resource_report(resource_report)
            else:
                logger.info("✅ 资源检查通过")

        except Exception as e:
            logger.error(f"资源检查失败: {e}")
            report["resource_check"] = {"error": str(e)}
            report["summary"]["blocking_issues"] += 1

        # 2. Prompt 质量审查
        visual_plan_json = self.project_dir / "visual_plan.json"
        if visual_plan_json.exists():
            logger.info("📝 审查 Visual Prompt 质量...")
            try:
                prompt_review_report = review_visual_plan(
                    visual_plan_json_path=str(visual_plan_json),
                    llm_client=self.client,
                    strategy="auto",  # 自动策略：轻检查 + 种子页深查 + 异常页必查
                    seed_pages=[1],  # 第一页作为种子页
                    content_plan_json_path=str(self.project_dir / "content_plan.json"),
                    style_config=None
                )
                report["prompt_review"] = prompt_review_report.to_dict()

                # 统计问题
                report["summary"]["total_issues"] += len(prompt_review_report.issues)
                report["summary"]["blocking_issues"] += prompt_review_report.stats.get("errors", 0)
                report["summary"]["warnings"] += prompt_review_report.stats.get("warnings", 0)

                if prompt_review_report.stats.get("errors", 0) > 0:
                    logger.warning(f"⚠️ Prompt 审查发现 {prompt_review_report.stats['errors']} 个错误")
                    print_review_report(prompt_review_report)
                else:
                    logger.info("✅ Prompt 审查通过")

            except Exception as e:
                logger.error(f"Prompt 审查失败: {e}")
                report["prompt_review"] = {"error": str(e)}
                report["summary"]["warnings"] += 1
        else:
            logger.warning("⚠️ visual_plan.json 不存在，跳过 Prompt 审查")
            report["prompt_review"] = {"skipped": "visual_plan.json not found"}

        # 3. Brief 一致性检查已移除（Brief 模块已删除）

        # 4. 判断是否通过
        report["passed"] = report["summary"]["blocking_issues"] == 0

        # 5. 打印汇总
        logger.info("\n" + "=" * 60)
        logger.info("最终意图审查报告")
        logger.info("=" * 60)
        if report["passed"]:
            logger.info("✅ 审查通过，可以执行")
        else:
            logger.error(f"❌ 审查未通过，发现 {report['summary']['blocking_issues']} 个阻塞问题")

        logger.info(f"\n问题统计:")
        logger.info(f"  总问题数: {report['summary']['total_issues']}")
        logger.info(f"  阻塞问题: {report['summary']['blocking_issues']}")
        logger.info(f"  警告: {report['summary']['warnings']}")
        logger.info("=" * 60 + "\n")

        return report

    def analyze_modification(self, modification_request: str, target_pages: Optional[List[int]] = None) -> Dict[str, Any]:
        """
        分析用户的修改请求，判断应该回退到哪个 Gate

        Args:
            modification_request: 用户的修改请求文本
            target_pages: 目标页面编号列表（如果用户指定了具体页面）

        Returns:
            分析结果，包含：
            - modification_type: 修改类型（content/visual/structure/image_assignment）
            - rollback_to_gate: 应回退到的 Gate（Content/Visual/Execute）
            - can_modify_single_page: 是否可以只修改单页
            - affected_pages: 受影响的页面编号列表
            - reason: 判断理由
            - suggestions: 操作建议
        """
        import json

        logger.info("🔍 分析修改请求")

        result = {
            "modification_type": None,
            "rollback_to_gate": None,
            "can_modify_single_page": False,
            "affected_pages": target_pages or [],
            "reason": "",
            "suggestions": []
        }

        # 构建分析 prompt
        prompt = f"""你是一个专业的产品经理，需要分析用户的修改请求，判断应该回退到哪个处理阶段。

系统有三个处理阶段（Gate）：
1. Content Gate: 内容规划阶段，生成 content_plan（页面结构、标题、文本内容）
2. Visual Gate: 视觉规划阶段，生成 visual_plan（图片选择、位置、视觉 prompt）
3. Execute Gate: 执行阶段，生成最终的 PPT 文件

用户修改请求：
{modification_request}

{"用户指定的目标页面: " + str(target_pages) if target_pages else "用户未指定具体页面"}

请分析这个修改请求，返回 JSON 格式的分析结果：

{{
  "modification_type": "修改类型，从以下选项中选择一个：content（内容修改）/ visual（视觉修改）/ structure（结构修改）/ image_assignment（图片归属修改）",
  "rollback_to_gate": "应回退到的 Gate，从以下选项中选择一个：Content / Visual / Execute",
  "can_modify_single_page": true/false,
  "reason": "判断理由（一句话说明为什么要回退到这个 Gate）",
  "suggestions": ["操作建议1", "操作建议2"]
}}

判断规则：
1. 修改类型判断：
   - content: 修改文字内容、标题、页面数量、页面顺序
   - visual: 修改图片选择、图片位置、视觉风格、prompt
   - structure: 修改整体结构、增删页面、重新组织内容
   - image_assignment: 修改图片与页面的对应关系、图片语义锚点

2. Gate 回退判断：
   - 回退到 Content: 如果修改涉及内容结构、页面数量、文字内容
   - 回退到 Visual: 如果修改只涉及图片选择、位置、视觉效果
   - 回退到 Execute: 如果修改只涉及 PPT 生成参数、格式调整

3. 单页修改判断：
   - 可以单页修改: 修改只影响特定页面的视觉效果或文字内容，不影响其他页面
   - 不能单页修改: 修改涉及整体结构、页面顺序、或影响多个页面的一致性

只返回 JSON，不要有其他内容。"""

        try:
            # 调用 LLM 分析
            response = self.client.chat.completions.create(
                model=DEFAULT_LLM_MODEL,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
            )

            result_text = response.choices[0].message.content.strip()

            # 提取 JSON
            if result_text.startswith("```"):
                result_text = result_text.split("```")[1]
                if result_text.startswith("json"):
                    result_text = result_text[4:]
            result_text = result_text.strip()

            analysis = json.loads(result_text)

            # 更新结果
            result["modification_type"] = analysis.get("modification_type")
            result["rollback_to_gate"] = analysis.get("rollback_to_gate")
            result["can_modify_single_page"] = analysis.get("can_modify_single_page", False)
            result["reason"] = analysis.get("reason", "")
            result["suggestions"] = analysis.get("suggestions", [])

            # 如果用户指定了页面，且判断可以单页修改，则只影响指定页面
            if target_pages and result["can_modify_single_page"]:
                result["affected_pages"] = target_pages
            else:
                # 否则影响所有页面
                result["affected_pages"] = self._get_all_page_numbers()

            logger.info(f"✅ 修改分析完成: {result['modification_type']} -> {result['rollback_to_gate']}")
            logger.info(f"   理由: {result['reason']}")
            logger.info(f"   单页修改: {result['can_modify_single_page']}")
            logger.info(f"   受影响页面: {result['affected_pages']}")

            return result

        except Exception as e:
            logger.error(f"修改分析失败: {e}")
            # 返回保守的默认值
            result["modification_type"] = "unknown"
            result["rollback_to_gate"] = "Content"
            result["can_modify_single_page"] = False
            result["reason"] = f"分析失败，保守起见回退到 Content Gate: {e}"
            result["suggestions"] = ["建议手动检查修改范围"]
            result["affected_pages"] = self._get_all_page_numbers()
            return result

    def _get_all_page_numbers(self) -> List[int]:
        """
        获取所有页面编号

        Returns:
            页面编号列表
        """
        import json

        # 尝试从 content_plan.json 获取
        content_plan_json = self.project_dir / "content_plan.json"
        if content_plan_json.exists():
            try:
                with open(content_plan_json, 'r', encoding='utf-8') as f:
                    content_plan = json.load(f)
                    slides = content_plan.get("slides", [])
                    return [slide.get("slide_number", i + 1) for i, slide in enumerate(slides)]
            except Exception:
                pass

        # 尝试从 visual_plan.json 获取
        visual_plan_json = self.project_dir / "visual_plan.json"
        if visual_plan_json.exists():
            try:
                with open(visual_plan_json, 'r', encoding='utf-8') as f:
                    visual_plan = json.load(f)
                    pages = visual_plan.get("pages", [])
                    return [page.get("page_number", i + 1) for i, page in enumerate(pages)]
            except Exception:
                pass

        # 默认返回空列表
        return []


# ──────────────────────────────────────────────
# ProjectState: 状态管理
# ──────────────────────────────────────────────

class ProjectState:
    """
    项目状态管理器。

    负责：
    - in-memory 状态（当前会话）
    - 原子性持久化到 _content_state.json
    - 快照 + 回滚机制
    """

    STATE_FILE = "_content_state.json"

    def __init__(self, project_dir: str):
        self.project_dir = Path(project_dir)
        self.project_dir.mkdir(parents=True, exist_ok=True)
        self._in_memory: Dict[str, Any] = {}
        self._dirty = False
        self._snapshots: List[Dict[str, str]] = []

        # 加载已有状态（如果存在）
        self._load()

    def get(self, key: str, default: Any = None) -> Any:
        """获取状态值（优先内存，其次持久化）"""
        if key in self._in_memory:
            return self._in_memory[key]

        # 尝试从持久化数据获取
        state_data = self._load()
        if state_data and "data" in state_data and key in state_data["data"]:
            return state_data["data"][key]

        return default

    def set(self, key: str, value: Any) -> None:
        """设置状态值（仅更新内存，标记 dirty）"""
        self._in_memory[key] = value
        self._dirty = True

    def get_all(self) -> Dict[str, Any]:
        """获取所有状态数据"""
        state_data = self._load()
        if state_data and "data" in state_data:
            return {**state_data["data"], **self._in_memory}
        return self._in_memory.copy()

    def persist(self) -> None:
        """原子性持久化到 _content_state.json"""
        if not self._dirty:
            return

        import tempfile
        import time

        # 构建完整状态
        state_data = self._load() or {
            "version": "2.0",
            "created_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            "snapshots": [],
        }

        # 更新时间戳
        state_data["updated_at"] = time.strftime("%Y-%m-%dT%H:%M:%S")
        state_data["data"] = {**state_data.get("data", {}), **self._in_memory}
        state_data["snapshots"] = self._snapshots

        # 原子性写入：先写临时文件，再替换
        tmp_path = self.project_dir / f"_state.tmp_{os.getpid()}.json"
        target_path = self.project_dir / self.STATE_FILE

        try:
            with open(tmp_path, 'w', encoding='utf-8') as f:
                json.dump(state_data, f, ensure_ascii=False, indent=2)

            # 原子性替换
            tmp_path.replace(target_path)
            self._dirty = False
        finally:
            # 清理临时文件
            if tmp_path.exists():
                tmp_path.unlink()

    def _load(self) -> Optional[Dict[str, Any]]:
        """从持久化文件加载状态"""
        state_file = self.project_dir / self.STATE_FILE
        if not state_file.exists():
            return None

        try:
            with open(state_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError):
            return None

    def snapshot(self, reason: str) -> Path:
        """创建版本快照（用于回滚）"""
        import time
        import shutil

        snapshot_id = f"v{int(time.time())}"
        snapshot_dir = self.project_dir / ".snapshots"
        snapshot_dir.mkdir(exist_ok=True)

        snapshot_path = snapshot_dir / f"{snapshot_id}_{reason}.json"

        # 复制当前状态到快照
        state_file = self.project_dir / self.STATE_FILE
        if state_file.exists():
            shutil.copy(str(state_file), str(snapshot_path))

        # 记录快照
        self._snapshots.append({
            "id": snapshot_id,
            "reason": reason,
            "path": str(snapshot_path),
            "timestamp": time.strftime("%Y-%m-%dT%H:%M:%S"),
        })

        self._dirty = True
        return snapshot_path

    def rollback(self, snapshot_path: Path) -> bool:
        """回滚到指定快照"""
        import shutil

        if not snapshot_path.exists():
            logger.error(f"快照文件不存在: {snapshot_path}")
            return False

        try:
            state_file = self.project_dir / self.STATE_FILE
            shutil.copy(str(snapshot_path), str(state_file))

            # 重新加载
            state_data = self._load()
            if state_data and "data" in state_data:
                self._in_memory = state_data["data"]
                self._snapshots = state_data.get("snapshots", [])

            self._dirty = False
            logger.info(f"✅ 已回滚到快照: {snapshot_path}")
            return True
        except Exception as e:
            logger.error(f"回滚失败: {e}")
            return False
