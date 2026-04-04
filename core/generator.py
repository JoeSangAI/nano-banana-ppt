"""
PPT 生成工作流 - 基于 Banana Slides 思路的简化版
支持从想法/大纲生成完整 PPT，并导出为 PPTX 格式
"""
import os
import json
import base64
import logging
import requests
from pathlib import Path
from typing import List, Dict, Optional, Union
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# LLM client is created in the calling code, not imported here

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def _fix_black_corners(img: Image.Image) -> Image.Image:
    """Fix black corners in generated images (placeholder function)"""
    return img


def _clamp(value: float, minimum: float, maximum: float) -> float:
    return max(minimum, min(value, maximum))


def _normalize_bbox(bbox: Dict) -> Optional[Dict[str, float]]:
    if not bbox:
        return None
    try:
        left = float(bbox.get("left", 0.0))
        top = float(bbox.get("top", 0.0))
        width = float(bbox.get("width", 0.0))
        height = float(bbox.get("height", 0.0))
    except (TypeError, ValueError):
        return None

    if width <= 0 or height <= 0:
        return None

    left = _clamp(left, 0.0, 1.0)
    top = _clamp(top, 0.0, 1.0)
    width = _clamp(width, 0.0, 1.0 - left)
    height = _clamp(height, 0.0, 1.0 - top)
    if width <= 0 or height <= 0:
        return None

    return {
        "left": round(left, 4),
        "top": round(top, 4),
        "width": round(width, 4),
        "height": round(height, 4),
    }


def _fit_bbox_within_region(candidate_bbox: Optional[Dict], allowed_bbox: Optional[Dict]) -> Optional[Dict[str, float]]:
    allowed = _normalize_bbox(allowed_bbox)
    if not allowed:
        return _normalize_bbox(candidate_bbox)

    candidate = _normalize_bbox(candidate_bbox) or allowed
    width = min(candidate["width"], allowed["width"])
    height = min(candidate["height"], allowed["height"])
    max_left = allowed["left"] + allowed["width"] - width
    max_top = allowed["top"] + allowed["height"] - height
    left = _clamp(candidate["left"], allowed["left"], max_left)
    top = _clamp(candidate["top"], allowed["top"], max_top)

    return {
        "left": round(left, 4),
        "top": round(top, 4),
        "width": round(width, 4),
        "height": round(height, 4),
    }


def _bbox_overlap_area(a: Optional[Dict], b: Optional[Dict]) -> float:
    box_a = _normalize_bbox(a)
    box_b = _normalize_bbox(b)
    if not box_a or not box_b:
        return 0.0

    ax2 = box_a["left"] + box_a["width"]
    ay2 = box_a["top"] + box_a["height"]
    bx2 = box_b["left"] + box_b["width"]
    by2 = box_b["top"] + box_b["height"]
    overlap_w = max(0.0, min(ax2, bx2) - max(box_a["left"], box_b["left"]))
    overlap_h = max(0.0, min(ay2, by2) - max(box_a["top"], box_b["top"]))
    return round(overlap_w * overlap_h, 6)


def _lock_overlay_bbox(
    original_image: Dict,
    calculated_image: Optional[Dict],
    blend_reserved_regions: List[Dict],
) -> Dict:
    merged_image = dict(original_image)
    candidate_bbox = None
    if calculated_image:
        candidate_bbox = calculated_image.get("dynamic_bounding_box") or calculated_image.get("bounding_box")

    allowed_bbox = (
        original_image.get("overlay_allowed_region")
        or original_image.get("bounding_box")
        or candidate_bbox
    )
    locked_bbox = _fit_bbox_within_region(candidate_bbox, allowed_bbox)

    if locked_bbox and any(_bbox_overlap_area(locked_bbox, reserved) > 0.0001 for reserved in blend_reserved_regions):
        fallback_bbox = _normalize_bbox(original_image.get("bounding_box")) or _normalize_bbox(original_image.get("overlay_allowed_region"))
        if fallback_bbox:
            locked_bbox = fallback_bbox

    if locked_bbox:
        merged_image["dynamic_bounding_box"] = locked_bbox
    return merged_image


def _merge_native_images_with_locked_regions(native_images: List[Dict], calculated_overlay_images: List[Dict]) -> List[Dict]:
    calculated_by_path = {
        img.get("path"): img for img in calculated_overlay_images if img.get("path")
    }
    blend_reserved_regions = [
        region
        for region in (
            _normalize_bbox(img.get("blend_reserved_region") or img.get("bounding_box"))
            for img in native_images
            if img.get("integration_mode", "overlay") == "blend"
        )
        if region
    ]

    merged_images = []
    for image in native_images:
        if image.get("integration_mode", "overlay") == "blend":
            merged_images.append(dict(image))
            continue
        calc_image = calculated_by_path.get(image.get("path"))
        merged_images.append(_lock_overlay_bbox(image, calc_image, blend_reserved_regions))
    return merged_images


class PPTGenerator:
    """PPT 生成器 - 核心工作流"""
    
    def __init__(self, api_key: str, api_base: str = None, slides_dir: str = "output/slides"):
        """
        Args:
            api_key: API 密钥（用于 LLM 调用）
            api_base: API Base URL（用于 LLM 调用）
            slides_dir: 临时幻灯片图片的保存目录
        """
        from openai import OpenAI
        self._api_key = api_key
        self.client = OpenAI(api_key=api_key, base_url=api_base) if api_base else OpenAI(api_key=api_key)
        self.text_model = "MiniMax-M2.7"
        self.visual_director_model = "MiniMax-M2.7"
        self.image_model = "gemini-3.1-flash-image-preview"  # 图片生成模型保持不变

        # 图片生成客户端（DeerAPI）
        image_gen_key = os.getenv("IMAGE_GEN_API_KEY") or api_key
        image_gen_base = os.getenv("IMAGE_GEN_API_BASE") or "https://api.deerapi.com/v1"
        self._image_gen_client = OpenAI(api_key=image_gen_key, base_url=image_gen_base)
        self._image_gen_api_key = image_gen_key

        self.output_dir = Path(slides_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def generate_image(self, description: str, aspect_ratio: str = "16:9", reference_images: List[Image.Image] = None, is_background_only: bool = False, resolution: str = "1K", native_images: List[Dict] = None) -> Image.Image:
        """
        生成单页PPT图片
        使用 Google REST API 直接调用（兼容 nano banana）
        resolution: "1K" | "2K" | "4K"，默认 1K
        """
        import requests

        resolution = (resolution or "1K").upper()
        if resolution not in ("1K", "2K", "4K"):
            resolution = "1K"

        logger.info(f"正在生成图片 ({resolution}): {description[:50]}...")
        if reference_images:
            logger.info(f"🎨 使用 {len(reference_images)} 张参考图片保持风格一致")

        # VisualAgent 已构造完整 prompt，此处仅追加分辨率等技术参数
        tech_suffix = f"\n\nTechnical: aspect ratio {aspect_ratio}, {resolution} resolution, sharp text rendering. CRITICAL: No black blocks, no solid black rectangles, seamless full-bleed composition."
        
        # Inject smart whitespace instructions based on native_images array
        if native_images and len(native_images) > 0:
            overlay_areas = []
            blend_areas = []
            for idx, img_conf in enumerate(native_images):
                layout = img_conf.get('layout')
                bbox = img_conf.get('bounding_box')
                integration_mode = img_conf.get('integration_mode', 'overlay')
                
                area_text = None
                if bbox:
                    # Translate bounding box to natural language roughly
                    left_pct = int(bbox.get('left', 0) * 100)
                    top_pct = int(bbox.get('top', 0) * 100)
                    w_pct = int(bbox.get('width', 0) * 100)
                    h_pct = int(bbox.get('height', 0) * 100)
                    
                    # Provide an even stronger spatial instruction
                    if left_pct > 50:
                        position = "the RIGHT SIDE"
                    elif left_pct + w_pct < 50:
                        position = "the LEFT SIDE"
                    else:
                        position = "the CENTER"
                        
                    area_text = f"a massive space on {position} (starting {left_pct}% from left, {top_pct}% from top, spanning {w_pct}% width)"
                elif layout:
                    # Legacy fallback
                    layout_prompts = {
                        "right_half": "the RIGHT SIDE of the image",
                        "left_half": "the LEFT SIDE of the image",
                        "center": "the CENTER area of the image",
                        "bottom_right": "the BOTTOM RIGHT corner"
                    }
                    if layout in layout_prompts:
                        area_text = layout_prompts[layout]
                
                if area_text:
                    if integration_mode == 'blend':
                        blend_areas.append(area_text)
                    else:
                        overlay_areas.append(area_text)
            
            if overlay_areas:
                areas_str = " and ".join(overlay_areas)
                tech_suffix += f" CRITICAL VISUAL CONSTRAINT: You ABSOLUTELY MUST leave {areas_str} completely BLANK and EMPTY. Do NOT generate ANY text, shapes, or complex backgrounds in this area. It must be a flat, solid color gradient because a photo will be pasted over it later."

            if blend_areas:
                areas_str = " and ".join(blend_areas)
                tech_suffix += f" CRITICAL INSTRUCTION: This image must contain ALL of the provided reference images (the ones in this prompt). Each reference image is a REAL PHOTOGRAPH that must appear COMPLETELY and UNCHANGED in the final output - do NOT modify, crop, redraw, distort, or alter any reference image. Keep their original content, colors, people, objects, text, and spatial composition exactly as-is. The reference images are the CONTENT of this slide - you are only generating the background environment around them. If a reference image shows a person, product, or scene, that exact person/product/scene must appear unchanged in your output."
            elif native_images and len(native_images) > 0:
                # 有原生图片但没有明确位置时，让 AI 根据图片内容自动决定最佳位置
                num_imgs = len(native_images)
                tech_suffix += f" CRITICAL INSTRUCTION: You have {num_imgs} reference images provided in this prompt. ALL {num_imgs} reference images must appear COMPLETELY and UNCHANGED in your output. Each one is a REAL PHOTOGRAPH - do NOT modify, redraw, blend, distort, or alter any of them. Their original content, colors, people, objects, text, and spatial composition must remain pixel-perfect. You are ONLY generating the background - the reference images are the content and must appear exactly as they are. Do not merge, composite, or regenerate the reference images in any way."


        full_prompt = description + tech_suffix

        # 使用 Google REST API 直接调用
        api_key = self.client.api_key if hasattr(self.client, 'api_key') else None
        if not api_key:
            api_key = getattr(self, '_api_key', None)

        if not api_key:
            raise ValueError("无法获取 API Key")

        is_openrouter = "openrouter" in str(getattr(self._image_gen_client, "base_url", "")).lower()
        is_deerapi = "deerapi" in str(getattr(self._image_gen_client, "base_url", "")).lower()

        if is_openrouter or is_deerapi:
            # 适配 OpenRouter 和 DeerAPI 的 OpenAI 格式调用
            # OpenRouter 需要 google/ 前缀，DeerAPI 直接使用模型名
            if is_openrouter:
                or_model = f"google/{self.image_model}" if not self.image_model.startswith("google/") else self.image_model
            else:  # DeerAPI
                or_model = self.image_model
            
            messages = [{"role": "user", "content": [{"type": "text", "text": full_prompt}]}]
            if reference_images:
                MAX_REFERENCE_IMAGES = 14  # Gemini 3 series via DeerAPI supports up to 14 reference images
                for ref_img in reference_images[:MAX_REFERENCE_IMAGES]:
                    buffered = BytesIO()
                    ref_img.save(buffered, format="PNG")
                    img_b64 = base64.b64encode(buffered.getvalue()).decode()
                    messages[0]["content"].append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{img_b64}"}
                    })
            
            import time
            for attempt in range(5):
                try:
                    resp = self._image_gen_client.chat.completions.create(
                        model=or_model,
                        messages=messages,
                        extra_body={"image_config": {"aspect_ratio": aspect_ratio, "imageSize": resolution}}
                    )
                    content = resp.choices[0].message.content
                    if content is None:
                        # 部分接口把图片放在 message.content 的 list 里
                        msg = resp.choices[0].message
                        if hasattr(msg, "content") and isinstance(getattr(msg, "content", None), list):
                            for part in msg.content:
                                if isinstance(part, dict):
                                    text = part.get("text") or part.get("type") == "image_url" and part.get("image_url", {}).get("url") or ""
                                    if text and "base64," in str(text):
                                        content = text
                                        break
                        if content is None:
                            raise ValueError("OpenRouter 未返回文本内容，可能当前模型不支持绘图或返回格式变更")
                    content = content or ""
                    # 提取 base64 图片数据
                    import re
                    match = re.search(r'data:image/[a-zA-Z]+;base64,([a-zA-Z0-9+/=]+)', content)
                    if match:
                        image_data = base64.b64decode(match.group(1))
                        return Image.open(BytesIO(image_data))
                    else:
                        raise ValueError("OpenRouter 返回的图像格式异常或无法解析 base64")
                except Exception as e:
                    logger.warning(f"OpenRouter 绘图失败 (尝试 {attempt+1}): {e}")
                    if attempt == 4: raise
                    time.sleep(2)
        else:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.image_model}:generateContent?key={api_key}"
            headers = {"Content-Type": "application/json"}
    
            # 构建 parts 数组：先文本，后参考图片
            parts = [{"text": full_prompt}]
    
            # 参考图片：[0]=模板风格图(可选) + [1..N]=原生融合图，Gemini API 最多支持 4 张
        MAX_REFERENCE_IMAGES = 14  # Gemini 3 series via DeerAPI supports up to 14 reference images
        if reference_images:
            for ref_img in reference_images[:MAX_REFERENCE_IMAGES]:
                buffered = BytesIO()
                ref_img.save(buffered, format="PNG")
                img_b64 = base64.b64encode(buffered.getvalue()).decode()
                parts.append({
                    "inlineData": {
                        "mimeType": "image/png",
                        "data": img_b64
                    }
                })

        # generationConfig: 支持 1K/2K/4K 分辨率
        generation_config = {
            "responseModalities": ["IMAGE"],
            "imageConfig": {
                "aspectRatio": aspect_ratio,
                "imageSize": resolution
            }
        }
        payload = {
            "contents": [{"parts": parts}],
            "generationConfig": generation_config,
        }
        
        # 添加重试机制（增加延迟，避免限流）
        max_retries = 5  # Increased from 3 to 5
        import time
        
        if hasattr(self, '_last_request_time'):
            elapsed = time.time() - self._last_request_time
            if elapsed < 1:
                time.sleep(1 - elapsed)
        
        for attempt in range(max_retries):
            try:
                response = requests.post(url, headers=headers, json=payload, timeout=180)
                
                if response.status_code == 200:
                    result = response.json()
                    candidates = result.get('candidates', [])
                    
                    if not candidates:
                        # Log the full result for debugging
                        logger.warning(f"Attempt {attempt + 1}: Empty candidates in response: {str(result)[:200]}...")
                        raise ValueError("生成失败：返回结果为空")
                    
                    content = candidates[0].get('content', {})
                    parts = content.get('parts', [])
                    
                    for part in parts:
                        # 兼容 REST API 可能返回的两种字段格式
                        img_b64 = None
                        if 'inlineData' in part:
                            img_b64 = part['inlineData']['data']
                        elif 'inline_data' in part:
                            img_b64 = part['inline_data']['data']
                        
                        if img_b64:
                            img_data = base64.b64decode(img_b64)
                            image = Image.open(BytesIO(img_data)).convert("RGB")
                            logger.info(f"图片生成成功: {image.size}")
                            self._last_request_time = time.time()
                            return image
                    
                    raise ValueError("未找到图片数据")
                else:
                    error_text = response.text
                    if attempt < max_retries - 1:
                        wait_time = 2 * (2 ** attempt)  # 2s, 4s, 8s, 16s...
                        logger.warning(f"API 请求失败 ({response.status_code})，重试 {attempt + 1}/{max_retries} (等待 {wait_time}s)...")
                        time.sleep(wait_time)
                        continue
                    raise Exception(f"API 请求失败 ({response.status_code}): {error_text}")
                    
            except (requests.exceptions.SSLError, requests.exceptions.ConnectionError, ValueError) as e:
                if attempt < max_retries - 1:
                    wait_time = 2 * (2 ** attempt)
                    logger.warning(f"请求异常 ({type(e).__name__})，重试 {attempt + 1}/{max_retries} (等待 {wait_time}s)...: {e}")
                    time.sleep(wait_time)
                    continue
                raise
                
        # 如果所有重试都失败
        raise Exception("图片生成失败：所有重试均失败")
    
    def upscale_image(self, image_path: str, resolution: str = "4K") -> bool:
        """
        使用 Gemini API 高保真放大已有图片。
        只放大不改变任何排版、文字、颜色或设计元素。
        返回是否成功。
        """
        import requests
        from PIL import Image
        import io
        import base64
        import time

        resolution = resolution.upper()
        if resolution not in ("2K", "4K"):
            logger.warning(f"⚠️ 分辨率参数错误 ({resolution})，不支持放大，保持原图。")
            return False

        if not os.path.exists(image_path):
            logger.error(f"❌ 找不到图片文件: {image_path}")
            return False

        logger.info(f"正在高保真放大图片至 {resolution}: {image_path}")

        prompt = (
            f"Upscale this image to {resolution} resolution. ACT AS A HIGH-FIDELITY UPSCALER. "
            "You must maintain all text, details, layouts, and colors exactly as they appear in the source image. "
            "Do NOT change any words, do NOT move any text, do NOT add or remove any design elements. "
            "Simply increase the resolution, sharpness, and clarity."
        )

        with open(image_path, "rb") as f:
            image_bytes = f.read()
        
        # 兼容现有的请求组装逻辑
        mime_type = "image/png"
        if str(image_path).lower().endswith(('.jpg', '.jpeg')):
            mime_type = "image/jpeg"
            
        b64_data = base64.b64encode(image_bytes).decode("utf-8")
        parts = [
            {"text": prompt},
            {"inlineData": {"mimeType": mime_type, "data": b64_data}}
        ]

        # generationConfig
        generation_config = {
            "responseModalities": ["IMAGE"],
            "imageConfig": {
                "aspectRatio": "16:9",
                "imageSize": resolution
            }
        }
        
        payload = {
            "contents": [{"parts": parts}],
            "generationConfig": generation_config,
        }

        api_key = self.client.api_key if hasattr(self.client, 'api_key') else getattr(self, '_api_key', None)
        if not api_key:
            logger.error("❌ 无法获取 API Key")
            return False

        api_base = self.client.base_url if hasattr(self.client, 'base_url') else "https://generativelanguage.googleapis.com/v1beta/openai"
        # 针对 Gemini API，直接构造 REST URL (非 OpenAI 兼容 URL)
        # 如果提供了 openai base url，提取主机名并重组
        if "googleapis.com" in str(api_base):
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.image_model}:generateContent?key={api_key}"
        else:
            # 对于第三方反代，假设它是直接反代的
            base = str(api_base).replace("/openai/v1", "").replace("/openai", "")
            url = f"{base}/models/{self.image_model}:generateContent?key={api_key}"

        headers = {"Content-Type": "application/json"}
        
        # 重试逻辑
        max_retries = 5
        base_wait = 3
        for attempt in range(max_retries):
            try:
                response = requests.post(url, headers=headers, json=payload, timeout=180)
                if response.status_code == 200:
                    data = response.json()
                    candidates = data.get("candidates", [])
                    if candidates:
                        parts = candidates[0].get("content", {}).get("parts", [])
                        for part in parts:
                            inline_data = part.get("inlineData") or part.get("inline_data")
                            if inline_data:
                                img_bytes = base64.b64decode(inline_data["data"])
                                img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                                # 覆盖保存原图
                                img.save(image_path)
                                logger.info(f"✅ 成功放大图片并覆盖保存: {image_path}")
                                return True
                    logger.error(f"❌ API返回异常数据格式: {str(data)[:200]}...")
                elif response.status_code == 429:
                    wait_time = base_wait * (2 ** attempt)
                    logger.warning(f"⚠️ API 速率限制 (429)，等待 {wait_time} 秒后重试...")
                    time.sleep(wait_time)
                    continue
                else:
                    logger.error(f"❌ API请求失败 ({response.status_code}): {response.text}")
                    if attempt < max_retries - 1:
                        wait_time = base_wait * (2 ** attempt)
                        time.sleep(wait_time)
                        continue
                    break
            except Exception as e:
                logger.error(f"❌ 图片放大生成出错: {e}")
                if attempt < max_retries - 1:
                    wait_time = base_wait * (2 ** attempt)
                    time.sleep(wait_time)
                
        return False

    def create_pptx(self, images: List[Image.Image], output_path: str) -> str:
        """
        Legacy: 将图片列表转换为 PPTX 文件 (全屏直出模式)
        """
        logger.info(f"正在创建PPTX文件 (Legacy): {output_path}")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9 比例
        
        for i, image in enumerate(images):
            # 保存临时图片
            temp_path = self.output_dir / f"temp_slide_{i}.png"
            image.save(temp_path, "PNG")
            
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加图片 (全屏铺满)
            slide.shapes.add_picture(
                str(temp_path),
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            # 清理临时文件
            temp_path.unlink()
        
        prs.save(output_path)
        logger.info(f"PPTX文件已保存: {output_path}")
        return output_path

    def _calculate_dynamic_layout(self, bg_img: Image.Image, native_images: List[Dict]) -> List[Dict]:
        """Use Vision LLM to calculate perfect bounding boxes based on the actual generated background."""
        import copy
        import base64
        import json
        from io import BytesIO
        
        # Deep copy to avoid mutating original plan dict globally if not needed
        updated_images = copy.deepcopy(native_images)
        
        try:
            buffered = BytesIO()
            bg_img.copy().convert("RGB").save(buffered, format="JPEG", quality=85)
            img_b64 = base64.b64encode(buffered.getvalue()).decode()
            
            img_details = []
            for idx, img_conf in enumerate(updated_images):
                path = img_conf.get('path')
                try:
                    from PIL import Image as PILImage
                    with PILImage.open(path) as n_img:
                        w, h = n_img.size
                        aspect = w / h
                    img_details.append(f"Image {idx+1}: Aspect Ratio {aspect:.2f} (Width/Height)")
                except Exception as e:
                    img_details.append(f"Image {idx+1}: Unknown aspect ratio (assume ~1.5)")

            prompt = f"""You are an expert presentation layout designer.
Look at the provided presentation slide image. This image already contains text and background graphics generated by another AI.
I need to place {len(updated_images)} native photo(s) on top of this slide. 

Your task is to find the ABSOLUTE PERFECT empty space (Safe Zone) on the slide to place the photo(s).
CRITICAL RULES:
1. Find the LARGEST possible rectangular empty space that DOES NOT OVERLAP ANY TEXT or important UI elements.
2. IMPORTANT: You must leave a generous padding/margin around your box. Do not let your box touch the text or the edges of the slide. If there is text on the left, start your box significantly to the right of it.
3. The bounding box must look visually balanced, aligning naturally with the text blocks (e.g., matching top/bottom margins).
4. Provide the coordinates as floats between 0.0 and 1.0 (where 0,0 is top-left, 1,1 is bottom-right).

Return ONLY a valid JSON array of bounding boxes in exactly the same order as the images listed above.
Example:
[
  {{"left": 0.55, "top": 0.20, "width": 0.40, "height": 0.60}}
]
"""
            logger.info(f"👁️ 正在使用视觉模型分析底层图片，寻找完美的排版位置...")

            # 适配 OpenRouter 前缀
            api_model = "gemini-3.1-pro-preview"
            if "openrouter" in str(self.client.base_url).lower() and not api_model.startswith("google/"):
                api_model = f"google/{api_model}"

            response = self.client.chat.completions.create(
                model=api_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{img_b64}"
                                }
                            }
                        ]
                    }
                ],
                temperature=0.1
            )
            
            content = response.choices[0].message.content.strip()
            start = content.find('[')
            end = content.rfind(']')
            if start != -1 and end != -1:
                json_str = content[start:end+1]
                boxes = json.loads(json_str)
                
                for i, box in enumerate(boxes):
                    if i < len(updated_images):
                        updated_images[i]['dynamic_bounding_box'] = box
                        logger.info(f"  🎯 视觉模型计算出图 {i+1} 完美坐标: {box}")
                        
        except Exception as e:
            logger.warning(f"视觉排版分析失败，将回退到原始排版: {e}")
            
        return updated_images

    def create_advanced_pptx(self, visual_plan: List[Dict], images: Dict[int, Image.Image], output_path: str, template_path: str = None, project_dir: str = None) -> str:
        """
        Advanced: 将视觉计划和图片组装为 PPTX
        - 支持 .pptx 模版母版映射
        - 支持原生文本填充 (Editable Text)
        - 支持 Hex 颜色强制
        """
        logger.info(f"正在创建高级 PPTX 文件: {output_path}")
        
        if template_path and template_path.lower().endswith('.pptx') and os.path.exists(template_path):
            logger.info(f"📦 使用 .pptx 模版: {template_path}")
            prs = Presentation(template_path)
        else:
            logger.info("📄 使用默认空白模版")
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
        # Helper: Find layout by name
        def get_layout(prs, layout_name_hints):
            for layout in prs.slide_layouts:
                name = layout.name.lower()
                for hint in layout_name_hints:
                    if hint.lower() in name:
                        return layout
            return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0] # Fallback
            
        # Helper: Hex to RGB
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

        for i, slide_plan in enumerate(visual_plan):
            page_type = slide_plan.get('type', 'content').lower()
            text_content = slide_plan.get('text_content', {})
            style_config = slide_plan.get('style_config', {})
            table_data = slide_plan.get('table_data') or text_content.get('table_data')
            visualization = slide_plan.get('visualization', '')

            # 对于模板页，使用空白布局，但稍后我们将添加文本框
            layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

            slide = prs.slides.add_slide(layout)

            page_num = slide_plan.get('page_num')
            img = images.get(page_num)

            if img:
                # 图表/普通页：添加背景图。优先使用已存在的 slide_XX.png（executor 已写入），
                # 避免重复写入 temp_slide 造成两套文件；仅当文件不存在时再保存（少见 edge case）
                slide_path = self.output_dir / f"slide_{page_num:02d}.png"
                if not slide_path.exists():
                    img.save(slide_path, "PNG")
                slide.shapes.add_picture(str(slide_path), 0, 0, prs.slide_width, prs.slide_height)

            # --- NEW: Add Editable Text Boxes for Template Slides ---
            if page_type.startswith('template_'):
                logger.info(f"✨ 为模板页 {page_num} 添加可编辑文本框和装饰...")

                def _clean_font_name(raw: str) -> str:
                    """Extract first usable font name from descriptive style library strings.
                    e.g. 'Tiempos Text / Copernicus / Georgia (for Headings ONLY)' → 'Georgia'
                    Prefers well-known cross-platform fonts when multiple options are listed.
                    Always returns the raw first token if none match the known list —
                    letting PPTX fall back gracefully on systems missing that font.
                    """
                    import re
                    KNOWN_FONTS = {
                        # Western
                        'arial', 'georgia', 'inter', 'helvetica', 'helvetica neue',
                        'times new roman', 'verdana', 'calibri', 'lato', 'montserrat',
                        'noto sans', 'noto serif', 'noto sans sc', 'noto serif sc',
                        'roboto', 'open sans', 'quicksand', 'nunito', 'cinzel',
                        'playfair display', 'cormorant garamond', 'dm sans', 'dm serif display',
                        'source sans pro', 'source serif pro', 'raleway', 'oswald',
                        # Chinese / CJK
                        'alibaba puhuiti', 'pingfang sc', 'pingfang tc', 'pingfang hk',
                        'hiragino sans gb', 'hiragino mincho pron', 'heiti sc', 'heiti tc',
                        'songti sc', 'kaiti sc', 'microsoft yahei', 'simsun', 'simhei',
                        'noto sans cjk sc', 'noto serif cjk sc', 'source han sans sc',
                        'source han serif sc', 'zcool qingke huangyou', 'zcool xiaowei',
                        'ma shan zheng', 'long cang', 'zhi mang xing',
                    }
                    parts = [p.strip() for p in raw.split('/')]
                    cleaned = []
                    for p in parts:
                        name = re.sub(r'\(.*?\)', '', p).strip()
                        if name:
                            cleaned.append(name)
                    # Prefer known/safe fonts from the list
                    for name in cleaned:
                        if name.lower() in KNOWN_FONTS:
                            return name
                    # No known font found — return the first candidate as-is
                    return cleaned[0] if cleaned else 'Arial'

                raw_fonts = style_config.get('fonts', ['Arial', 'Arial'])
                title_font = _clean_font_name(raw_fonts[0]) if raw_fonts else 'Arial'
                body_font = _clean_font_name(raw_fonts[-1]) if raw_fonts else 'Arial'
                palette = style_config.get('palette', ['#FFFFFF', '#000000', '#CCCCCC', '#666666'])

                bg_color_hex = palette[0] if palette else '#FFFFFF'

                def _hex_to_rgb_tuple(hex_col):
                    h = hex_col.lstrip('#')
                    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

                def _blend_hex(hex_a, hex_b, ratio=0.7):
                    """Blend hex_a toward hex_b by ratio (0.7 = 70% a, 30% b)."""
                    ra, ga, ba = _hex_to_rgb_tuple(hex_a)
                    rb, gb, bb = _hex_to_rgb_tuple(hex_b)
                    r = int(ra * ratio + rb * (1 - ratio))
                    g = int(ga * ratio + gb * (1 - ratio))
                    b = int(ba * ratio + bb * (1 - ratio))
                    return f'#{r:02X}{g:02X}{b:02X}'

                def _luminance(hex_col):
                    r, g, b = _hex_to_rgb_tuple(hex_col)
                    return 0.2126 * r + 0.7152 * g + 0.0722 * b

                def _saturation(hex_col):
                    """Simple HSL-based saturation (0-255 scale)."""
                    r, g, b = [x / 255.0 for x in _hex_to_rgb_tuple(hex_col)]
                    cmax, cmin = max(r, g, b), min(r, g, b)
                    return (cmax - cmin) * 255

                bg_lum = _luminance(bg_color_hex)

                # --- 4-Level Color Hierarchy (aligned with AI-rendered pages) ---
                # palette convention: [0]=bg, [1]=primary text (sent to Gemini as "Primary Text"),
                # [2+]=accent/secondary colors.
                # Use palette[1] directly as title color to match AI pages exactly.
                text_color_hex = palette[1] if len(palette) > 1 else '#222222'
                if abs(bg_lum - _luminance(text_color_hex)) < 60:
                    text_color_hex = '#222222' if bg_lum > 128 else '#FFFFFF'

                # Accent: pick the most saturated non-bg non-text palette color
                best_accent_hex = None
                best_accent_score = -1
                for c in palette:
                    if c.lower() == bg_color_hex.lower() or c.lower() == text_color_hex.lower():
                        continue
                    sat = _saturation(c)
                    contrast = abs(bg_lum - _luminance(c))
                    score = sat * 0.6 + contrast * 0.4
                    if score > best_accent_score:
                        best_accent_score = score
                        best_accent_hex = c
                if best_accent_hex is None:
                    best_accent_hex = palette[2] if len(palette) > 2 else '#666666'
                accent_color_hex = best_accent_hex

                # L1 Title: primary text color (matches AI pages)
                text_color = hex_to_rgb(text_color_hex)
                # L2 Subtitle: title blended slightly toward bg (90:10)
                subtitle_color_hex = _blend_hex(text_color_hex, bg_color_hex, ratio=0.88)
                subtitle_color = hex_to_rgb(subtitle_color_hex)
                # L3 Body: softer blend (72:28)
                body_color_hex = _blend_hex(text_color_hex, bg_color_hex, ratio=0.72)
                body_text_color = hex_to_rgb(body_color_hex)
                # L4 Muted/placeholder: even lighter (55:45)
                muted_color_hex = _blend_hex(text_color_hex, bg_color_hex, ratio=0.55)
                muted_color = hex_to_rgb(muted_color_hex)
                # Accent color for decorative elements, tags, bullet markers
                accent_color = hex_to_rgb(accent_color_hex)
                # Tag label text: ensure contrast against accent background
                tag_text_color = hex_to_rgb('#FFFFFF') if _luminance(accent_color_hex) < 150 else hex_to_rgb('#222222')

                logger.info(f"模板字体: 标题={title_font}, 正文={body_font}")
                logger.info(f"模板配色: 标题={text_color_hex}, 正文={body_color_hex}, 装饰={accent_color_hex}")
                
                # Determine alignment based on layout hint
                layout_hint = slide_plan.get('layout', 'centered_headline')
                is_centered = 'center' in layout_hint
                
                from pptx.enum.text import PP_ALIGN
                title_align = PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT
                
                if page_type == 'template_content':
                    # Add Top Decorative Line (More subtle)
                    line = slide.shapes.add_shape(
                        9, # MSO_SHAPE.LINE
                        Inches(1),
                        Inches(0.4),
                        Inches(14),
                        Inches(0)
                    )
                    line.line.color.rgb = accent_color
                    line.line.width = Pt(1.5)

                    # 明确的引导标签 (Template Instruction Label)
                    tag_box = slide.shapes.add_shape(1, Inches(1), Inches(0.55), Inches(1.5), Inches(0.3))
                    tag_box.fill.solid()
                    tag_box.fill.fore_color.rgb = accent_color
                    tag_box.line.color.rgb = accent_color
                    tf = tag_box.text_frame
                    tf.text = "自由编辑页"
                    p = tf.paragraphs[0]
                    p.font.name = body_font
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = tag_text_color
                    p.alignment = PP_ALIGN.CENTER
                    from pptx.enum.text import MSO_ANCHOR
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                    # Add Title Placeholder
                    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(14), Inches(1.2))
                    tf = title_box.text_frame
                    tf.text = "单击此处输入您的标题"
                    p = tf.paragraphs[0]
                    p.font.name = title_font
                    p.font.size = Pt(48)
                    p.font.bold = True
                    p.font.color.rgb = text_color
                    p.alignment = PP_ALIGN.LEFT

                    # Add Subtitle hint (L2)
                    sub_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(14), Inches(0.5))
                    tf = sub_box.text_frame
                    tf.text = "单击此处添加正文内容。您可以自由发挥。"
                    p = tf.paragraphs[0]
                    p.font.name = body_font
                    p.font.size = Pt(22)
                    p.font.color.rgb = subtitle_color

                    # Add Body Placeholder with accent bullet markers (L3 + L4)
                    body_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(14), Inches(4.8))
                    tf = body_box.text_frame
                    tf.word_wrap = True
                    body_lines = [
                        ("• 按下回车键可以添加新的段落", body_text_color),
                        ("• 使用这一页来补充 AI 没有覆盖到的重要细节", body_text_color),
                        ("• 所有的字体和颜色已经预设为您当前的主题风格", body_text_color),
                    ]
                    tf.text = ""
                    for li, (line_text, line_color) in enumerate(body_lines):
                        if li == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.font.name = body_font
                        p.font.size = Pt(20)
                        p.space_after = Pt(14)
                        # Accent-colored bullet marker + body-colored text
                        from pptx.oxml.ns import qn
                        run_bullet = p.add_run()
                        run_bullet.text = "•  "
                        run_bullet.font.color.rgb = accent_color
                        run_bullet.font.size = Pt(20)
                        run_bullet.font.bold = True
                        run_body = p.add_run()
                        run_body.text = line_text.lstrip("• ")
                        run_body.font.color.rgb = line_color
                        run_body.font.size = Pt(20)
                        run_body.font.name = body_font

                elif page_type == 'template_split':
                    # Add Top Decorative Line
                    line = slide.shapes.add_shape(9, Inches(1), Inches(0.4), Inches(14), Inches(0))
                    line.line.color.rgb = accent_color
                    line.line.width = Pt(1.5)

                    # 明确的引导标签 (Template Instruction Label)
                    tag_box = slide.shapes.add_shape(1, Inches(1), Inches(0.55), Inches(1.8), Inches(0.3))
                    tag_box.fill.solid()
                    tag_box.fill.fore_color.rgb = accent_color
                    tag_box.line.color.rgb = accent_color
                    tf = tag_box.text_frame
                    tf.text = "图文分栏页 (可编辑)"
                    p = tf.paragraphs[0]
                    p.font.name = body_font
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = tag_text_color
                    p.alignment = PP_ALIGN.CENTER
                    from pptx.enum.text import MSO_ANCHOR
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                    # Add Title Placeholder
                    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(14), Inches(1.2))
                    tf = title_box.text_frame
                    tf.text = "此处输入分栏排版标题"
                    p = tf.paragraphs[0]
                    p.font.name = title_font
                    p.font.size = Pt(48)
                    p.font.bold = True
                    p.font.color.rgb = text_color
                    p.alignment = PP_ALIGN.LEFT

                    # Left Content Box — subtitle hint (L2) + body with accent bullets (L3)
                    left_sub = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(6.5), Inches(0.5))
                    tf = left_sub.text_frame
                    tf.text = "此处添加正文说明"
                    p = tf.paragraphs[0]
                    p.font.name = body_font
                    p.font.size = Pt(22)
                    p.font.color.rgb = subtitle_color

                    left_box = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(6.5), Inches(4.7))
                    tf = left_box.text_frame
                    tf.word_wrap = True
                    left_lines = [
                        "您可以在左侧输入核心观点或数据描述，右侧插入相关图片或图表作为呼应。",
                        "这种排版非常适合进行概念解释",
                        "也可以用于对比说明",
                    ]
                    tf.text = ""
                    for li, line_text in enumerate(left_lines):
                        if li == 0:
                            p = tf.paragraphs[0]
                            p.font.name = body_font
                            p.font.size = Pt(18)
                            p.font.color.rgb = body_text_color
                            p.space_after = Pt(16)
                            run = p.runs[0] if p.runs else p.add_run()
                            run.text = line_text
                        else:
                            p = tf.add_paragraph()
                            p.space_after = Pt(12)
                            run_b = p.add_run()
                            run_b.text = "•  "
                            run_b.font.color.rgb = accent_color
                            run_b.font.size = Pt(18)
                            run_b.font.bold = True
                            run_t = p.add_run()
                            run_t.text = line_text
                            run_t.font.color.rgb = body_text_color
                            run_t.font.size = Pt(18)
                            run_t.font.name = body_font

                    # Right Image Placeholder — muted color (L4)
                    right_box = slide.shapes.add_shape(1, Inches(8.5), Inches(2.5), Inches(6.5), Inches(5.5))
                    right_box.fill.solid()
                    placeholder_bg = _blend_hex(bg_color_hex, text_color_hex, ratio=0.92)
                    right_box.fill.fore_color.rgb = hex_to_rgb(placeholder_bg)
                    right_box.line.color.rgb = accent_color
                    right_box.line.width = Pt(1)
                    right_box.line.dash_style = 4
                    
                    tf = right_box.text_frame
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf.text = ""
                    p_icon = tf.paragraphs[0]
                    p_icon.alignment = PP_ALIGN.CENTER
                    run = p_icon.add_run()
                    run.text = "🖼️"
                    run.font.size = Pt(48)
                    p_hint = tf.add_paragraph()
                    p_hint.alignment = PP_ALIGN.CENTER
                    p_hint.space_before = Pt(12)
                    run = p_hint.add_run()
                    run.text = "删除此框并在此处插入您的图片"
                    run.font.name = body_font
                    run.font.size = Pt(16)
                    run.font.bold = True
                    run.font.color.rgb = muted_color
                    p_sub = tf.add_paragraph()
                    p_sub.alignment = PP_ALIGN.CENTER
                    run = p_sub.add_run()
                    run.text = "(Insert > Picture)"
                    run.font.name = body_font
                    run.font.size = Pt(14)
                    run.font.color.rgb = muted_color

            # Add Logo as a separate, movable PPTX shape (not burned into image)
            logo_path = slide_plan.get('logo_path')
            import os
            if logo_path and os.path.exists(logo_path) and page_type != 'background_only':
                logo_loc = (slide_plan.get('logo_location') or 'Top-Right').lower()
                
                # Define bounding box for logo
                max_logo_w = Inches(1.8)
                max_logo_h = Inches(0.65)
                margin_x = Inches(0.5)
                margin_y = Inches(0.4)
                
                try:
                    from PIL import Image as PILImage
                    _logo = PILImage.open(logo_path)
                    logo_aspect = _logo.width / _logo.height
                    
                    # Calculate dimensions preserving aspect ratio within bounding box
                    # First try setting width to max
                    calc_w = max_logo_w
                    calc_h = calc_w / logo_aspect
                    
                    # If height exceeds max_height, scale down based on height instead
                    if calc_h > max_logo_h:
                        calc_h = max_logo_h
                        calc_w = calc_h * logo_aspect
                        
                    logo_w = calc_w
                    logo_h = calc_h
                except Exception as e:
                    logger.warning(f"Failed to read logo for aspect ratio: {e}")
                    logo_w = Inches(1.2)
                    logo_h = Inches(0.45)

                if 'left' in logo_loc and 'top' in logo_loc:
                    lx, ly = margin_x, margin_y
                elif 'right' in logo_loc and 'bottom' in logo_loc:
                    lx = prs.slide_width - logo_w - margin_x
                    ly = prs.slide_height - logo_h - margin_y
                elif 'left' in logo_loc and 'bottom' in logo_loc:
                    lx, ly = margin_x, prs.slide_height - logo_h - margin_y
                else: # Default Top-Right
                    lx = prs.slide_width - logo_w - margin_x
                    ly = margin_y

                slide.shapes.add_picture(logo_path, lx, ly, logo_w, logo_h)

            # Add speaker notes logic moved up
            speaker_notes = slide_plan.get('speaker_notes')
            if speaker_notes:
                try:
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = speaker_notes
                except Exception as e:
                    logger.warning(f"无法添加演讲者备注到第 {i+1} 页: {e}")

            # --- NEW: Add Multiple Native Images ---
            native_images = slide_plan.get('native_images', [])
            if not native_images and slide_plan.get('native_image'):
                native_images = [slide_plan.get('native_image')]

            # Only overlay images need placement. Blend images are already baked into the background.
            overlay_images = [img_conf for img_conf in native_images if img_conf.get("integration_mode", "overlay") == "overlay"]
            if overlay_images and img:
                calculated_overlay_images = self._calculate_dynamic_layout(img, overlay_images)
                native_images = _merge_native_images_with_locked_regions(native_images, calculated_overlay_images)

            for img_conf in native_images:
                if img_conf.get("integration_mode", "overlay") == "blend":
                    continue
                img_path = img_conf.get('path')
                
                # Check if it's an http path that hasn't been resolved yet
                if img_path and img_path.startswith("http"):
                    import os
                    from urllib.parse import urlparse
                    
                    # Try to find it relative to the original content file
                    content_file = slide_plan.get("style_config", {}).get("_meta", {}).get("content_file", "")
                    if not content_file:
                        # Sometimes passed globally via visual_plan root metadata
                        pass
                        
                    # Let's see if we can find it in the same directory as the script execution or project dir
                    filename = os.path.basename(urlparse(img_path).path)
                    
                    # Search up to find the document directory
                    possible_paths = [
                        filename,
                        os.path.join(project_dir, filename) if project_dir else filename,
                    ]
                    
                    # Add original document dir if available from the plan
                    # We can't access meta directly here easily, but we can search common locations
                    import glob
                    found = False
                    for search_path in [f"**/{filename}"]:
                        matches = glob.glob(search_path, recursive=True)
                        if matches:
                            img_path = matches[0]
                            found = True
                            break
                            
                    if not found:
                        logger.warning(f"无法找到网络图片对应的本地文件: {img_path}")
                        continue
                
                if not img_path or not os.path.exists(img_path):
                    logger.warning(f"原生图片不存在或路径为空: {img_path}")
                    continue
                    
                layout = img_conf.get('layout', 'center')
                bbox = img_conf.get('bounding_box')
                
                try:
                    from PIL import Image as PILImage
                    import tempfile
                    
                    # Convert WebP to PNG temporarily if needed since python-pptx doesn't support WebP
                    is_temp_file = False
                    if img_path.lower().endswith('.webp'):
                        try:
                            webp_img = PILImage.open(img_path)
                            temp_fd, temp_path = tempfile.mkstemp(suffix='.png')
                            os.close(temp_fd)
                            webp_img.save(temp_path, format="PNG")
                            img_path = temp_path
                            is_temp_file = True
                        except Exception as e:
                            logger.warning(f"无法转换 WebP 图片 {img_path}: {e}")
                            continue
                            
                    native_img = PILImage.open(img_path)
                    img_w, img_h = native_img.size
                    aspect = img_w / img_h
                    
                    sw = prs.slide_width
                    sh = prs.slide_height
                    
                    margin = Inches(0.5)
                    
                    # 1. Resolve target bounding box (left, top, max_width, max_height)
                    dynamic_bbox = img_conf.get('dynamic_bounding_box')
                    active_bbox = dynamic_bbox if dynamic_bbox else bbox
                    
                    if active_bbox:
                        # Semantic coordinate system (0.0 - 1.0)
                        target_l = sw * active_bbox.get('left', 0)
                        target_t = sh * active_bbox.get('top', 0)
                        max_w = sw * active_bbox.get('width', 1.0)
                        max_h = sh * active_bbox.get('height', 1.0)
                        logger.info(f"    计算原生图片坐标: target_l={target_l}, target_t={target_t}, max_w={max_w}, max_h={max_h} (来自 VLM: {bool(dynamic_bbox)})")
                    else:
                        # Legacy enum system
                        if layout == 'right_half':
                            box = (sw / 2 + margin/2, margin, sw / 2 - margin*1.5, sh - margin*2)
                        elif layout == 'left_half':
                            box = (margin, margin, sw / 2 - margin*1.5, sh - margin*2)
                        elif layout == 'bottom_right':
                            box = (sw * 0.6, sh * 0.5, sw * 0.4 - margin, sh * 0.5 - margin)
                        elif layout == 'fullscreen':
                            box = (0, 0, sw, sh)
                        else: # center
                            box = (margin*2, margin*2, sw - margin*4, sh - margin*4)
                        target_l, target_t, max_w, max_h = box
                        
                    # 2. Calculate fitted dimensions preserving aspect ratio
                    if max_h == 0: max_h = 1 # prevent div by zero
                    target_aspect = max_w / max_h
                    if aspect > target_aspect:
                        # Image is wider than target box, so width is the limiting factor
                        final_w = max_w
                        final_h = max_w / aspect
                    else:
                        # Image is taller than target box, so height is the limiting factor
                        final_h = max_h
                        final_w = max_h * aspect
                        
                    # 3. Align within the target box
                    if dynamic_bbox:
                        # VLM calculates the safe zone. We MUST perfectly center the image inside this safe zone.
                        # This prevents the image from sticking to the edge and maintains the VLM's intended margins.
                        final_l = target_l + (max_w - final_w) / 2
                        final_t = target_t + (max_h - final_h) / 2
                    else:
                        # Old mechanical logic fallback
                        if active_bbox:
                            left_pct = active_bbox.get('left', 0)
                            if left_pct < 0.2:
                                # Align left
                                final_l = target_l
                            elif left_pct > 0.5:
                                # Align right
                                final_l = target_l + (max_w - final_w)
                            else:
                                # Center horizontal
                                final_l = target_l + (max_w - final_w) / 2
                        else:
                            final_l = target_l + (max_w - final_w) / 2
                            
                        # Always center vertically for now
                        final_t = target_t + (max_h - final_h) / 2
                        
                    slide.shapes.add_picture(img_path, final_l, final_t, final_w, final_h)
                    logger.info(f"  已插入多图排版图片: {active_bbox if active_bbox else layout}")
                    
                    # Clean up temp file if we created one
                    if is_temp_file:
                        try:
                            os.remove(img_path)
                        except:
                            pass
                            
                except Exception as e:
                    logger.warning(f"无法插入多图排版图片 {img_path}: {e}")
                    if 'is_temp_file' in locals() and is_temp_file:
                        try:
                            os.remove(img_path)
                        except:
                            pass

        prs.save(output_path)
        logger.info(f"高级 PPTX 文件已保存: {output_path}")
        return output_path

if __name__ == "__main__":
    pass
