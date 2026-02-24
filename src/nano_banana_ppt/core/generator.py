"""
PPT 生成工作流 - 基于 Banana Slides 思路的简化版
支持从想法/大纲生成完整 PPT，并导出为 PPTX 格式
"""
import os
import json
import base64
import logging
import requests
from pathlib import Path
from typing import List, Dict, Optional
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from openai import OpenAI

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTGenerator:
    """PPT 生成器 - 核心工作流"""
    
    def __init__(self, api_key: str, api_base: str = None):
        """
        初始化 PPT 生成器
        
        Args:
            api_key: API 密钥
            api_base: API Base URL（可选，用于代理，文本生成时使用）
        """
        self._api_key = api_key  # 保存 API Key 用于图片生成
        # 文本生成仍使用 OpenAI SDK（兼容 OpenAI 格式）
        self.client = OpenAI(
            api_key=api_key,
            base_url=api_base or "https://generativelanguage.googleapis.com/v1beta/openai",
            timeout=120.0,
            max_retries=3
        )
        self.text_model = "gemini-3-flash-preview"  # 用于大纲生成等轻量任务
        self.visual_director_model = "gemini-3-pro-preview"  # 用于视觉转译，需要更强的推理能力
        self.image_model = "gemini-3-pro-image-preview"
        self.output_dir = Path("output")
        self.output_dir.mkdir(exist_ok=True)
    
    def generate_outline(self, idea: str, language: str = "zh") -> List[Dict]:
        """
        从想法生成大纲
        
        Args:
            idea: 一句话想法
            language: 语言代码 (zh/en/ja)
            
        Returns:
            大纲列表，每个元素包含 title 和 description
        """
        prompt = f"""请根据以下想法，生成一个专业的PPT大纲。

想法：{idea}

要求：
1. 生成8-12页PPT的大纲
2. 每页包含：标题和简要描述
3. 结构清晰，逻辑连贯
4. 使用{language}输出

请以JSON格式输出，格式如下：
[
  {{"title": "第1页标题", "description": "第1页内容描述"}},
  {{"title": "第2页标题", "description": "第2页内容描述"}},
  ...
]
"""
        
        logger.info("正在生成大纲...")
        response = self.client.chat.completions.create(
            model=self.text_model,
            messages=[
                {"role": "system", "content": "你是一个专业的PPT内容策划专家。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7
        )
        
        content = response.choices[0].message.content
        logger.info(f"大纲生成完成: {content[:100]}...")
        
        # 解析 JSON
        try:
            # 尝试提取 JSON 部分
            json_match = content.find('[')
            if json_match != -1:
                json_str = content[json_match:]
                json_end = json_str.rfind(']') + 1
                if json_end > 0:
                    json_str = json_str[:json_end]
                    outline = json.loads(json_str)
                    return outline
        except Exception as e:
            logger.warning(f"JSON解析失败，尝试手动解析: {e}")
        
        # 如果 JSON 解析失败，尝试从文本中提取
        outline = []
        lines = content.split('\n')
        for line in lines:
            if 'title' in line.lower() or '标题' in line:
                # 简单提取逻辑
                pass
        
        return outline if outline else [{"title": "封面", "description": idea}]
    
    def identify_page_type(self, page: Dict) -> str:
        """
        识别页面类型：封面、目录、内容页
        
        Args:
            page: 页面信息字典
            
        Returns:
            "cover" | "toc" | "content"
        """
        title = page.get('title', '').lower()
        description = page.get('description', '').lower()
        
        # 封面识别
        if '封面' in title or 'cover' in title or page.get('page_num') == 1:
            return "cover"
        
        # 目录识别
        if '目录' in title or 'toc' in title or 'table of contents' in title or 'contents' in title:
            return "toc"
        
        # 默认内容页
        return "content"
    
    def extract_style_from_image(self, image: Image.Image, page_type: str = "cover") -> str:
        """
        从图片中提取风格关键词（用于风格一致性）
        
        Args:
            image: PIL Image 对象
            page_type: 页面类型 ("cover" | "toc" | "content")
            
        Returns:
            风格描述字符串
        """
        # 将图片转换为 base64
        buffered = BytesIO()
        image.save(buffered, format="PNG")
        img_b64 = base64.b64encode(buffered.getvalue()).decode()
        
        if page_type == "cover":
            prompt = f"""分析这张PPT封面图片的视觉风格，提取**核心设计语言**（配色、字体风格、设计调性）。

请提取以下维度的风格特征：
1. **配色方案**：主色调、辅助色、背景色
2. **设计风格**：Minimalist / Corporate / Tech / Creative 等
3. **整体调性**：High-end / Modern / Professional 等

**重要**：提取的风格要适合应用到内容页，避免过于"封面化"的元素（如大标题、全屏背景等）。
提取核心的配色和设计语言即可。

请用简洁的英文关键词描述，格式如：
"Deep blue background, golden accents, minimalist design, high-end corporate style, clean typography"

只输出风格描述，不要其他解释。"""
        else:
            prompt = f"""分析这张PPT页面的视觉风格，提取关键风格元素。

请提取以下维度的风格特征：
1. **配色方案**：主色调、辅助色、背景色
2. **设计风格**：Minimalist / Corporate / Tech / Creative 等
3. **视觉元素**：图标风格、图表类型、字体感觉
4. **整体调性**：High-end / Modern / Professional 等

请用简洁的英文关键词描述，格式如：
"Deep blue background, golden accents, minimalist design, high-end corporate style, clean typography"

只输出风格描述，不要其他解释。"""

        try:
            # 使用 Gemini 的多模态能力分析图片
            response = self.client.chat.completions.create(
                model=self.visual_director_model,
                messages=[
                    {"role": "system", "content": "你是一个专业的视觉风格分析师。"},
                    {"role": "user", "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}}
                    ]}
                ],
                temperature=0.3
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            logger.warning(f"风格提取失败，使用默认风格: {e}")
            return "Professional corporate style, deep blue background, golden accents, minimalist"
    
    def adapt_style_for_content(self, cover_style: str) -> str:
        """
        将封面风格适配为内容页风格
        
        Args:
            cover_style: 封面风格描述
            
        Returns:
            适配后的内容页风格描述
        """
        prompt = f"""将以下封面风格适配为适合内容页的风格。

封面风格：{cover_style}

要求：
1. **保留核心元素**：配色方案、设计调性、字体风格
2. **移除封面元素**：大标题、全屏背景、过于装饰性的元素
3. **适配内容页**：适合信息密度高、可读性强的内容页
4. **保持一致性**：确保与封面风格在视觉上统一

请输出适配后的风格描述，只输出风格关键词，不要解释。"""

        try:
            response = self.client.chat.completions.create(
                model=self.visual_director_model,
                messages=[
                    {"role": "system", "content": "你是一个专业的PPT风格适配专家。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.5
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            logger.warning(f"风格适配失败，使用原风格: {e}")
            return cover_style
    
    def translate_to_visual_descriptor(self, page_content: Dict, style_constraint: str = None) -> str:
        """
        [Visual Director] 将页面内容转译为视觉描述符 (Bilingual Mode)
        
        Args:
            page_content: 页面内容字典
            style_constraint: 风格约束
            
        Returns:
            str: 包含中文文本指令和英文视觉描述的 Prompt
        """
        # 1. 提取并清洗关键文本 (中文)
        title = page_content.get('title', '').replace('封面', '').replace('目录', '').replace('P1', '').replace('P2', '').strip()
        description = page_content.get('description', '')
        # 移除 "核心内容：" 等前缀
        if "：" in description:
            description = description.split("：", 1)[1]
            
        key_points = page_content.get('content', [])[:3]
        clean_points = []
        for p in key_points:
            # 清洗列表项中的元数据前缀
            clean_p = p.split('：', 1)[1] if '：' in p else p
            clean_points.append(clean_p)
        
        # 2. 构建中文文本指令块
        chinese_text_instruction = f"""
【必须渲染的中文内容 (MUST RENDER CHINESE TEXT)】
请确保以下中文文字清晰、准确地出现在画面中：
- 大标题 (Title): "{title}"
- 副标题/金句 (Subtitle): "{description}"
"""
        if clean_points:
            chinese_text_instruction += "正文列表 (Body Text):\n" + "\n".join([f'- "{p}"' for p in clean_points])

        # 3. 构建 Prompt 生成指令
        prompt = f"""You are a world-class Presentation Designer using Google's Nano Banana Pro model.
Your task is to write a **Natural Language Prompt** to generate a stunning PPT slide.

【Content Context】
Topic: {title}
Core Message: {description}

【Visual Direction】
1. **Natural Language Description**: Describe the slide as if briefing a human designer. Use complete sentences.
   - Example: "Create a professional business slide featuring a deep blue gradient background with golden geometric accents."
   
2. **Style Injection**:
   - {style_constraint if style_constraint else "Professional, High-end, Corporate, Minimalist, Deep Navy & Gold palette"}
   - Layout: Modern grid system, clear hierarchy, balanced whitespace.
   - Typography: Clean sans-serif for body, elegant serif for titles.

3. **Negative Constraints (CRITICAL)**:
   - Do NOT render labels like "封面", "P1", "目录", "核心内容", "Cover", "Page 1".
   - Do NOT render placeholder text like "Lorem Ipsum".
   - Do NOT translate the Chinese text into English.

【Output Format】
Directly output the final Image Generation Prompt.
Structure:
[Visual Description in English] ...
[Text Rendering Instruction]
{chinese_text_instruction}

Start with: "A professional presentation slide..."
"""

        try:
            response = self.client.chat.completions.create(
                model=self.visual_director_model,
                messages=[
                    {"role": "system", "content": "You are an expert Prompt Engineer for Nano Banana Pro. You are fluent in both English and Chinese."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            logger.error(f"视觉转译失败: {e}")
            return f"A professional presentation slide titled '{title}', text: '{description}'. {style_constraint or 'Business style'}"
    
    def generate_page_description(self, title: str, description: str, context: str = "") -> str:
        """
        生成单页详细描述（用于图片生成）
        
        Args:
            title: 页面标题
            description: 页面描述
            context: 上下文信息（可选）
            
        Returns:
            详细的页面描述文本
        """
        prompt = f"""请为以下PPT页面生成详细的视觉描述，用于AI图片生成。

标题：{title}
内容：{description}
{f'上下文：{context}' if context else ''}

要求：
1. 描述要具体、详细，包含视觉元素（颜色、布局、图标等）
2. 适合16:9比例的PPT页面
3. 包含所有需要在页面上展示的文字内容
4. 描述要专业、美观

请直接输出描述文本，不要包含其他说明。"""
        
        response = self.client.chat.completions.create(
            model=self.text_model,
            messages=[
                {"role": "system", "content": "你是一个专业的PPT视觉设计师。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.8
        )
        
        return response.choices[0].message.content
    
    def generate_image(self, description: str, aspect_ratio: str = "16:9", reference_images: List[Image.Image] = None, is_background_only: bool = False) -> Image.Image:
        """
        生成单页PPT图片
        使用 Google REST API 直接调用（兼容 nano banana）
        
        Args:
            description: 图片描述prompt
            aspect_ratio: 宽高比
            reference_images: 参考图片列表（用于风格一致性）
            is_background_only: 是否为纯背景页（无文字、无图标）
        """
        import requests
        
        logger.info(f"正在生成图片: {description[:50]}...")
        if reference_images:
            logger.info(f"🎨 使用 {len(reference_images)} 张参考图片保持风格一致")
        
        # 构建图片生成 prompt
        if is_background_only:
            # 纯背景页：使用简洁的prompt，不添加文字和图标相关要求
            full_prompt = f"""{description}

要求：
- 这是一个纯背景纹理图片
- 比例：{aspect_ratio}
- 分辨率：2K (2560x1440) High Resolution
- 高质量纹理，无文字，无图标，无前景对象
- 4K high quality, detailed texture, aspect ratio {aspect_ratio}."""
        else:
            # 普通内容页：使用完整的prompt要求
            full_prompt = f"""{description}

要求：
- 这是一个专业的商业提案PPT页面
- 比例：{aspect_ratio}
- 分辨率：2K (2560x1440) High Resolution
- 风格：现代、专业、商务、高端
- 配色：建议使用深色背景（如深蓝、黑色）配金色/白色文字，体现高端感
- 布局：标题清晰醒目，内容层次分明
- 所有文字内容必须清晰可读，字体大小适中
- 可以适当添加图标、图表等视觉元素
- 整体视觉效果要符合"战略合作方案"的商务调性

请生成完整的PPT页面图片，包含所有文字内容。4K high quality, detailed, sharp text rendering. Aspect ratio {aspect_ratio}."""

        
        # 使用 Google REST API 直接调用
        api_key = self.client.api_key if hasattr(self.client, 'api_key') else None
        if not api_key:
            # 从 base_url 中提取或使用默认
            api_key = getattr(self, '_api_key', None)
        
        if not api_key:
            raise ValueError("无法获取 API Key")
        
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.image_model}:generateContent?key={api_key}"
        headers = {"Content-Type": "application/json"}
        
        # 构建 parts 数组：先文本，后参考图片
        parts = [{"text": full_prompt}]
        
        # 添加参考图片（如果提供）- 用于风格一致性
        if reference_images:
            for ref_img in reference_images[:1]:  # 限制为最多1张，避免干扰文字
                buffered = BytesIO()
                ref_img.save(buffered, format="PNG")
                img_b64 = base64.b64encode(buffered.getvalue()).decode()
                parts.append({
                    "inlineData": {
                        "mimeType": "image/png",
                        "data": img_b64
                    }
                })
        
        payload = {
            "contents": [{"parts": parts}],
            "generationConfig": {
                "responseModalities": ["IMAGE"],
                # 尝试请求更高分辨率或更好的图像质量参数（如果 API 支持）
                # 注意：Gemini API 的具体参数可能随版本变化，目前主要靠 prompt 控制
            }
        }
        
        # 添加重试机制（增加延迟，避免限流）
        max_retries = 3
        import time
        
        # 在请求前添加短暂延迟，避免触发 API 限流
        if hasattr(self, '_last_request_time'):
            elapsed = time.time() - self._last_request_time
            if elapsed < 2:  # 至少间隔2秒
                time.sleep(2 - elapsed)
        
        for attempt in range(max_retries):
            try:
                response = requests.post(url, headers=headers, json=payload, timeout=180)
                
                if response.status_code == 200:
                    result = response.json()
                    candidates = result.get('candidates', [])
                    
                    if not candidates:
                        raise ValueError("生成失败：返回结果为空")
                    
                    content = candidates[0].get('content', {})
                    parts = content.get('parts', [])
                    
                    for part in parts:
                        # 兼容 REST API 可能返回的两种字段格式
                        img_b64 = None
                        if 'inlineData' in part:
                            img_b64 = part['inlineData']['data']
                        elif 'inline_data' in part:
                            img_b64 = part['inline_data']['data']
                        
                        if img_b64:
                            img_data = base64.b64decode(img_b64)
                            image = Image.open(BytesIO(img_data))
                            logger.info(f"图片生成成功: {image.size}")
                            # 记录成功请求时间
                            self._last_request_time = time.time()
                            return image
                    
                    raise ValueError("未找到图片数据")
                else:
                    error_text = response.text
                    if attempt < max_retries - 1:
                        logger.warning(f"API 请求失败 ({response.status_code})，重试 {attempt + 1}/{max_retries}...")
                        import time
                        time.sleep(2 ** attempt)  # 指数退避
                        continue
                    raise Exception(f"API 请求失败 ({response.status_code}): {error_text}")
                    
            except (requests.exceptions.SSLError, requests.exceptions.ConnectionError) as e:
                if attempt < max_retries - 1:
                    logger.warning(f"网络错误，重试 {attempt + 1}/{max_retries}...: {e}")
                    import time
                    time.sleep(2 ** attempt)  # 指数退避
                    continue
                raise
            except Exception as e:
                if attempt < max_retries - 1:
                    logger.warning(f"生成失败，重试 {attempt + 1}/{max_retries}...: {e}")
                    import time
                    time.sleep(2 ** attempt)
                    continue
                raise
                
        # 如果所有重试都失败
        raise Exception("图片生成失败：所有重试均失败")
    
    def create_pptx(self, images: List[Image.Image], output_path: str) -> str:
        """
        将图片列表转换为 PPTX 文件 (全屏直出模式)
        
        Args:
            images: PIL Image 对象列表
            output_path: 输出文件路径
            
        Returns:
            输出文件路径
        """
        logger.info(f"正在创建PPTX文件: {output_path}")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9 比例
        
        for i, image in enumerate(images):
            # 保存临时图片
            temp_path = self.output_dir / f"temp_slide_{i}.png"
            image.save(temp_path, "PNG")
            
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加图片 (全屏铺满)
            slide.shapes.add_picture(
                str(temp_path),
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            # 清理临时文件
            temp_path.unlink()
        
        prs.save(output_path)
        logger.info(f"PPTX文件已保存: {output_path}")
        return output_path
        """
        将图片列表转换为 PPTX 文件
        
        Args:
            images: PIL Image 对象列表
            output_path: 输出文件路径
            
        Returns:
            输出文件路径
        """
        logger.info(f"正在创建PPTX文件: {output_path}")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9 比例
        
        for i, image in enumerate(images):
            # 保存临时图片
            temp_path = self.output_dir / f"temp_slide_{i}.png"
            image.save(temp_path, "PNG")
            
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加图片
            slide.shapes.add_picture(
                str(temp_path),
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            # 清理临时文件
            temp_path.unlink()
        
        prs.save(output_path)
        logger.info(f"PPTX文件已保存: {output_path}")
        return output_path
    
    def generate_ppt(self, idea: str, output_name: str = None, language: str = "zh") -> str:
        """
        完整工作流：从想法生成完整PPT
        
        Args:
            idea: 一句话想法
            output_name: 输出文件名（不含扩展名）
            language: 语言代码
            
        Returns:
            生成的PPTX文件路径
        """
        logger.info(f"开始生成PPT: {idea}")
        
        # 1. 生成大纲
        outline = self.generate_outline(idea, language)
        logger.info(f"大纲生成完成，共 {len(outline)} 页")
        
        # 2. 生成每页图片
        images = []
        for i, page in enumerate(outline):
            logger.info(f"正在生成第 {i+1}/{len(outline)} 页: {page.get('title', '未知')}")
            
            # 生成详细描述
            page_description = self.generate_page_description(
                page.get('title', ''),
                page.get('description', ''),
                context=f"这是第{i+1}页，共{len(outline)}页"
            )
            
            # 生成图片
            try:
                image = self.generate_image(page_description)
                images.append(image)
            except Exception as e:
                logger.error(f"第 {i+1} 页生成失败: {e}")
                # 创建一个占位图
                placeholder = Image.new('RGB', (1920, 1080), color='lightgray')
                images.append(placeholder)
        
        # 3. 生成PPTX
        if not output_name:
            output_name = f"PPT_{idea[:20].replace(' ', '_')}"
        
        output_path = self.output_dir / f"{output_name}.pptx"
        self.create_pptx(images, str(output_path))
        
        logger.info(f"PPT生成完成: {output_path}")
        return str(output_path)


def main():
    """示例用法"""
    import sys
    
    # 从环境变量或命令行参数获取 API Key
    api_key = os.getenv("OPENAI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    api_base = os.getenv("OPENAI_API_BASE") or "https://generativelanguage.googleapis.com/v1beta/openai"
    
    if not api_key:
        print("错误: 请设置 OPENAI_API_KEY 或 GOOGLE_API_KEY 环境变量")
        print("或者通过命令行参数传入: python ppt_generator.py '你的想法' [API_KEY]")
        sys.exit(1)
    
    # 获取想法
    if len(sys.argv) > 1:
        idea = sys.argv[1]
    else:
        idea = input("请输入你的PPT想法: ")
    
    # 如果提供了第二个参数，作为 API Key
    if len(sys.argv) > 2:
        api_key = sys.argv[2]
    
    # 创建生成器并生成PPT
    generator = PPTGenerator(api_key, api_base)
    output_path = generator.generate_ppt(idea)
    
    print(f"\n✅ PPT生成成功！")
    print(f"📁 文件位置: {output_path}")


if __name__ == "__main__":
    main()

