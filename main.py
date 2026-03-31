"""
Auto Pipeline (RFC v2.2)
全自动 PPT 生成入口：支持 plan / execute 两阶段调用，兼容 Agent 和交互式终端

输出目录: ~/Desktop/AI output/ppt/{date}_{project_name}/
  {date}_{project_name}.pptx   ← 最终交付的 PPTX
  content_plan.md              ← 内容大纲
  master_plan.md               ← 视觉计划
  plan.json                    ← 技术执行计划
  template_assets/             ← 模板提取物
  slides/                      ← 生成的页面图片
"""
import os
import sys
import json
import logging
import re
from pathlib import Path

try:
    from dotenv import load_dotenv
except ImportError:
    load_dotenv = None


def _find_and_load_env() -> bool:
    """
    智能查找并加载 .env 文件
    优先级: 工具目录 → 上级目录 → 工作区根 → nano-banana-ppt skill 目录
    """
    if load_dotenv is None:
        return False
    tool_dir = Path(__file__).resolve().parent
    root = Path(os.getcwd())
    env_locations = [
        tool_dir / ".env",
        root / ".env",
        Path(__file__).resolve().parents[2] / ".env",  # workspace root
        Path.home() / ".cursor" / "skills" / "nano-banana-ppt" / ".env",
        Path.home() / ".claude" / "skills" / "ppt-generator" / ".env",
    ]
    for env_path in env_locations:
        if env_path.exists():
            load_dotenv(dotenv_path=env_path, override=True)
            logger = logging.getLogger(__name__)
            logger.info(f"Loaded .env from: {env_path}")
            return True
    load_dotenv(override=True)
    return False


# ── Auto-bootstrap: 确保无论从哪个目录调用，模块路径都正确 ──────────────────────
# 解析 skill 的实际位置（跟随 symlink 到真实路径）
_actual_skill_dir = Path(__file__).resolve().parent  # e.g. /Users/Joe_1/Desktop/nano-banana-ppt
_cwd = Path.cwd()

# 判断当前是否在一个有效工作区（有 tools/ 目录）
_cwd_has_tools = (_cwd / "tools").is_dir() or (_cwd / "tools" / "nano_banana_ppt").is_symlink()

# 判断当前是否就在 skill 目录内
_in_skill_dir = _cwd == _actual_skill_dir or _cwd.resolve() == _actual_skill_dir

# 如果在 skill 目录内，或当前工作区没有 tools/ 结构 → 自动切换到 skill 所在目录
if _in_skill_dir or not _cwd_has_tools:
    # 切换到 skill 的父目录（期望是 Desktop 或其他有 tools/ 的工作区）
    _potential_workspace = _actual_skill_dir.parent  # e.g. Desktop
    _workspace_has_tools = (_potential_workspace / "tools").is_dir()

    if _workspace_has_tools:
        os.chdir(_potential_workspace)
        sys.path.insert(0, str(_potential_workspace.resolve()))
    else:
        # 兜底：在 skill 目录本身创建 tools/ 结构
        _tools_link = _actual_skill_dir / "tools" / "nano_banana_ppt"
        _tools_link.parent.mkdir(parents=True, exist_ok=True)
        if not _tools_link.exists():
            _tools_link.symlink_to(_actual_skill_dir)
        os.chdir(_actual_skill_dir)
        sys.path.insert(0, str(_actual_skill_dir.resolve()))

# ── 环境变量加载（在 bootstrap 之后）───────────────────────────────────────────
_find_and_load_env()
from tools.nano_banana_ppt.agents.narrative import NarrativeAgent
from tools.nano_banana_ppt.agents.visual import VisualAgent
from tools.nano_banana_ppt.agents.template import TemplateAgent
from tools.nano_banana_ppt.core.executor import execute_plan
from tools.nano_banana_ppt.utils.llm_client import reset_session
from tools.nano_banana_ppt.utils.review_plan import (
    build_master_plan_from_content_plan,
    build_content_review_md,
    parse_review_md,
    derive_technical_plan,
    REVIEW_MD_FILENAME,
)
from tools.nano_banana_ppt.utils.content_validator import validate_content_visual_consistency

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# 固定的输出目录
DEFAULT_OUTPUT_BASE = Path.home() / "Desktop" / "AI output" / "ppt"


def _resolve_project_dir(content_file: str, output_name: str = None) -> tuple[str, Path]:
    """根据内容文件或指定名称，生成项目名和项目目录"""
    from datetime import date
    date_prefix = date.today().strftime("%Y%m%d")

    # 固定的输出目录
    OUTPUT_BASE = Path.home() / "Desktop" / "AI output" / "ppt"

    # 如果用户指定了 output_name，直接使用
    if output_name:
        name = output_name
    else:
        # 从 content_file 路径提取名称
        content_path = Path(content_file)

        # 检查是否已经在 output 目录中（用户可能在已有项目目录下运行）
        # 例如: /Users/Joe_1/Desktop/AI output/ppt/20260327_内容电商大逃杀/content_plan.md
        # 或: /Users/Joe_1/Desktop/AI output/ppt/20260327_20260327_内容电商大逃杀时代/content_plan.md
        parent_parts = content_path.parts
        for i, part in enumerate(parent_parts):
            # 跳过日期前缀的目录（如 20260327_内容电商大逃杀 或 20260327_20260327_xxx）
            if i > 0 and parent_parts[i-1] == "ppt" and part.startswith("202"):
                # 这是一个日期命名的项目目录，提取主题名称
                # 格式可能是: 20260327_主题名 或 20260327_20260327_主题名
                name_part = "_".join(parent_parts[i:]).replace("_content_plan", "").replace("_content", "")
                # 去除开头的日期前缀
                name_part = re.sub(r'^\d{8}_', '', name_part)
                name_part = re.sub(r'^\d{8}_\d{8}_', '', name_part)
                if name_part:
                    name = name_part
                    break
        else:
            # 没有匹配到日期格式的目录，使用文件名（去掉扩展名和 content_plan 后缀）
            name = content_path.stem
            name = re.sub(r'_content_plan$', '', name)
            name = re.sub(r'_content$', '', name)

    # 如果名称以日期开头，避免重复
    if re.match(r'^\d{8}_', name):
        dir_name = name
    else:
        dir_name = f"{date_prefix}_{name}"

    project_dir = OUTPUT_BASE / dir_name
    project_dir.mkdir(parents=True, exist_ok=True)
    return name, project_dir


# ──────────────────────────────────────────────
#  Phase 1: Plan
# ──────────────────────────────────────────────

def generate_content_plan(content_file: str, template_file: str = None,
                          logo_file: str = None, output_name: str = None, page_count: int = None,
                          briefing: str = None, reuse_existing: bool = False):
    """
    Phase 1.1: 内容规划阶段。
    生成 content_plan.md，仅供审阅大纲内容，不涉及视觉风格。
    确认后再运行 plan-visual 生成视觉计划。
    """
    reset_session()
    project_name, project_dir = _resolve_project_dir(content_file, output_name)
    tpl_assets_dir = project_dir / "template_assets"
    tpl_assets_dir.mkdir(exist_ok=True)
    content_md_path = project_dir / "content_plan.md"

    print(f"\n🚀 Nano Banana 2 — Phase 1.1: Content Plan")
    print(f"📂 项目名称: {project_name}")
    print(f"📂 项目目录: {project_dir}")
    print(f"📄 内容源: {content_file}")
    print(f"🎨 模版源: {template_file or '无 (AI 自动设计)'}")
    print(f"🏷️ Logo 源: {logo_file or '未指定'}")
    if briefing:
        print(f"📝 用户意图 (Briefing): {briefing[:60]}{'...' if len(briefing) > 60 else ''}")

    api_key = os.getenv("OPENAI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    api_base = os.getenv("OPENAI_API_BASE") or "https://generativelanguage.googleapis.com/v1beta/openai"
    if not api_key:
        print("❌ 错误: 请设置 OPENAI_API_KEY 或 GOOGLE_API_KEY 环境变量")
        return None

    narrative_agent = NarrativeAgent(api_key, api_base, project_dir=str(project_dir))
    template_agent = TemplateAgent(api_key, api_base, output_dir=str(tpl_assets_dir))

    if not os.path.exists(content_file):
        print(f"❌ 错误: 内容文件不存在 {content_file}")
        return None

    # 支持 PDF：用 PyMuPDF 提取文本；否则按 UTF-8 文本读取
    if content_file.lower().endswith('.pdf'):
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(content_file)
            parts = []
            for page in doc:
                parts.append(page.get_text())
            num_pages = len(doc)
            doc.close()
            content_context = "\n\n".join(parts)
            if not content_context.strip():
                print("❌ 错误: PDF 中未提取到文本，请检查文件或改用已导出的 .md/.txt")
                return None
            logger.info(f"已从 PDF 提取 {num_pages} 页文本")
        except Exception as e:
            logger.error(f"PDF 文本提取失败: {e}")
            print(f"❌ 错误: PDF 读取失败 ({e})，请确认已安装 pymupdf: pip install pymupdf")
            return None
    else:
        with open(content_file, 'r', encoding='utf-8') as f:
            content_context = f.read()

    # Step 1 & 2: 并行执行内容分析和模版解析（两者互不依赖）
    template_info = None
    assets = {}

    from concurrent.futures import ThreadPoolExecutor, as_completed

    def _analyze():
        print("\n🔍 [Step 1/3] 正在分析文档内容...")
        result = narrative_agent.analyze_content(content_context)
        logger.info(f"自动推导参数: {json.dumps(result, ensure_ascii=False)}")
        return result

    def _parse_template():
        if template_file and os.path.exists(template_file):
            print(f"\n🖼️ [Step 2/3] 正在解析模版文件...")
            try:
                info = template_agent.process_template(template_file)
                logger.info("模版解析成功")
                return info
            except Exception as e:
                logger.error(f"模版解析失败: {e}")
        elif template_file:
            logger.warning(f"模版文件不存在: {template_file}，将使用 AI 自动设计")
        return None

    with ThreadPoolExecutor(max_workers=2) as pool:
        fut_analyze = pool.submit(_analyze)
        fut_template = pool.submit(_parse_template)
        inferred = fut_analyze.result()
        template_info = fut_template.result()

    if template_info:
        assets['template_file'] = template_file
        assets['logo_path'] = template_info.get('logo_path')
        assets['template_images'] = template_info.get('reference_images')

    if logo_file and os.path.exists(logo_file):
        assets['logo_path'] = logo_file
        print(f"🏷️ 使用用户指定 Logo: {logo_file} (优先于模版截取)")
    elif logo_file:
        logger.warning(f"Logo 文件不存在: {logo_file}，将回退到模版截取")

    # If we have a logo but no template, extract brand colors
    brand_colors = []
    if assets.get('logo_path') and not template_info:
        from tools.nano_banana_ppt.utils.image_utils import extract_dominant_colors
        print("\n🎨 正在从 Logo 提取品牌色...")
        brand_colors = extract_dominant_colors(assets['logo_path'], num_colors=3)
        if brand_colors:
            print(f"   提取到的品牌色: {', '.join(brand_colors)}")
            inferred['brand_colors'] = brand_colors

    # 构建约束
    constraints = {
        "target_audience": inferred.get("target_audience", "通用受众"),
        "presentation_type": inferred.get("presentation_type", "商业演示"),
        "duration": inferred.get("duration", "15分钟"),
        "style_preference": inferred.get("style_preference", "专业严谨"),
        "page_count": str(page_count) if page_count else "10",
        "briefing": briefing,
    }

    # Step 3: 生成叙事大纲
    print("\n📝 [Step 3/3] 正在构建叙事架构...")

    # 预检测：如果检测到完整大纲且未指定 --reuse-existing，提示用户
    if not reuse_existing:
        complete_plan_info = NarrativeAgent.detect_complete_content_plan(content_context)
        if complete_plan_info["is_complete"] and complete_plan_info["confidence"] >= 0.8:
            print(f"""
✅ 检测到完整大纲（{complete_plan_info['page_count']}页）
   - 包含标题/正文/备注: {'是' if complete_plan_info['has_speaker_notes'] else '否'}
   - 包含数据表格: {'是' if complete_plan_info['has_tables'] else '否'}
   - 置信度: {complete_plan_info['confidence']:.1%}

💡 建议使用 --reuse-existing 参数直接复用，避免内容被修改：
   python3 -m tools.nano_banana_ppt.main plan-content "{content_file}" --reuse-existing

⏳ 继续执行将使用"解析模式"（可能会修改内容）...
            """)
            import time
            time.sleep(2)  # 给用户2秒时间看到提示

    narrative_outline = narrative_agent.generate_narrative_outline(content_context, constraints, content_file_path=content_file, reuse_existing=reuse_existing)
    print(f"✅ 叙事大纲生成完成，共 {len(narrative_outline)} 页")

    meta = {
        "project_name": project_name,
        "project_dir": str(project_dir),
        "content_file": content_file,
        "template_file": assets.get('template_file'),
        "logo_file": assets.get('logo_path'),
    }

    # 保存 _content_state.json
    state = {
        "inferred": inferred,
        "assets": assets,
        "meta": meta,
        "template_info": template_info,
        "narrative_outline": narrative_outline,
        "constraints": constraints
    }
    state_file = project_dir / "_content_state.json"
    with open(state_file, "w", encoding='utf-8') as f:
        json.dump(state, f, ensure_ascii=False, indent=2)

    # 保存 content_plan.md
    content_md_content = build_content_review_md(narrative_outline, meta)
    with open(content_md_path, "w", encoding='utf-8') as f:
        f.write(content_md_content)

    print(f"\n{'='*60}")
    print(f"📋 内容草稿出来了，您看看文字和插图对不对？")
    print(f"   已保存到: {content_md_path}")
    print(f"{'='*60}")
    for p in narrative_outline:
        p_num = p['page_num']
        p_type = p.get('type', 'content')
        tc = p.get('text_content', {})
        title = tc.get('headline') or p.get('title', '')
        print(f"  P{p_num} [{p_type.upper():8s}] {title}")
    print(f"{'='*60}")
    print(f"\n⛔ STOP — 请审阅 content_plan.md 确认无误后再继续！")
    print(f"   下一步: python main.py plan-visual \"{project_dir}\" [--style xxx]")

    return str(content_md_path)


def generate_visual_plan(plan_dir: str, style_preference: str = None):
    """
    Phase 1.5: 视觉规划阶段。
    读取内容计划，生成设计系统和视觉主张，输出 master_plan.md。
    """
    project_dir = Path(plan_dir)
    state_file = project_dir / "_content_state.json"
    review_md_path = project_dir / REVIEW_MD_FILENAME  # master_plan.md

    # 如果 _content_state.json 不存在，尝试从 content_plan.md 重建
    if not state_file.exists():
        content_md_path = project_dir / "content_plan.md"
        if content_md_path.exists():
            print(f"📄 检测到 content_plan.md 但无状态文件，自动重建 _content_state.json...")
            with open(content_md_path, 'r', encoding='utf-8') as f:
                md_text = f.read()
            parsed = parse_review_md(md_text)
            if parsed and parsed.get("pages"):
                # 从 content_plan.md 提取主题名作为项目名
                project_name = project_dir.name
                # 去掉日期前缀
                project_name = re.sub(r'^\d{8}_', '', project_name)
                project_name = re.sub(r'^\d{8}_\d{8}_', '', project_name)
                # 构建最小化状态
                state = {
                    "inferred": {"project_name": project_name},
                    "assets": {},
                    "meta": {"source_file": str(content_md_path)},
                    "template_info": None,
                    "narrative_outline": parsed.get("pages", []),
                    "constraints": {}
                }
                with open(state_file, 'w', encoding='utf-8') as f:
                    json.dump(state, f, ensure_ascii=False, indent=2)
                print(f"✅ 状态文件已创建，继续执行...")
            else:
                print(f"❌ 错误: 找不到内容状态文件，也无法从 content_plan.md 解析。")
                return None
        else:
            print(f"❌ 错误: 找不到内容状态文件 {state_file}。请先运行 plan-content。")
            return None

    with open(state_file, 'r', encoding='utf-8') as f:
        state = json.load(f)

    inferred = state.get("inferred", {})
    assets = state.get("assets", {})
    meta = state.get("meta", {})
    template_info = state.get("template_info")
    narrative_outline = state.get("narrative_outline", [])
    constraints = state.get("constraints", {})

    # 尝试从用户可能编辑过的 content_plan.md 读取最新大纲（覆盖 state 中的旧数据）
    content_md_path = project_dir / "content_plan.md"
    if content_md_path.exists():
        with open(content_md_path, 'r', encoding='utf-8') as f:
            md_text = f.read()
        parsed = parse_review_md(md_text)
        if parsed and parsed.get("pages"):
            narrative_outline = parsed.get("pages")
            print(f"📄 已读取最新编辑的 content_plan.md 大纲")

    print(f"\n🚀 Nano Banana 2 — Phase 1.5: Visual Plan")
    print(f"📂 项目目录: {project_dir}")
    print(f"🎨 风格偏好: {style_preference or 'AI 自动推导'}")

    api_key = os.getenv("OPENAI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    api_base = os.getenv("OPENAI_API_BASE") or "https://generativelanguage.googleapis.com/v1beta/openai"
    if not api_key:
        print("❌ 错误: 请设置 OPENAI_API_KEY 或 GOOGLE_API_KEY 环境变量")
        return None

    # 支持 prompt_mode 参数 (verbose 或 minimal)
    prompt_mode = os.getenv("PROMPT_MODE", "verbose")
    visual_agent = VisualAgent(api_key, api_base, prompt_mode=prompt_mode)

    # Step 1: 生成风格定义
    print("\n🎨 [Step 1/3] 正在定义视觉风格 (AI Minting)...")
    constraints["style_preference"] = style_preference or inferred.get("style_preference", "专业商务")
    brand_colors = inferred.get("brand_colors", [])
    if brand_colors:
        constraints["brand_colors"] = brand_colors

    style_definition = visual_agent.define_style(constraints, assets, template_info)
    if isinstance(style_definition, tuple):
        style_desc_str, style_config = style_definition
    else:
        style_desc_str = str(style_definition)
        style_config = {"description": style_desc_str, "palette": [], "mode": "ai_minting"}

    # Step 2: 为每页生成具体的「配图/画面」描述（Visual Director 提案）
    print("\n📋 [Step 2/4] 正在生成每页视觉提案 (Visual Director)...\n⏳ 这可能需要 1-2 分钟...")
    from tools.nano_banana_ppt.utils.review_plan import generate_per_slide_visual_suggestions
    per_slide_suggestions = generate_per_slide_visual_suggestions(
        narrative_outline, style_config, api_key, api_base
    )
    if per_slide_suggestions:
        print(f"✅ 已生成 {len(per_slide_suggestions)} 页的视觉描述")
        for pnum, suggestion in sorted(per_slide_suggestions.items()):
            print(f"   P{pnum}: {suggestion[:60]}...")
    else:
        print("⚠️ 未生成视觉描述，将留空供用户补充")

    # Step 3: 生成视觉主张 (Art Director Manifesto) — 仅供用户审阅
    print("\n📋 [Step 3/4] 正在生成视觉主张 (Art Director Manifesto)...")
    from tools.nano_banana_ppt.utils.review_plan import generate_design_manifesto

    is_template_mode = bool(meta.get("template_file"))
    parsed_stub = {"pages": narrative_outline, "style": style_config}
    manifesto_dict = generate_design_manifesto(
        parsed=parsed_stub,
        template_mode=is_template_mode,
        client=visual_agent.client
    )
    manifesto_text = manifesto_dict.get("chinese_proposal", "")

    # 写回 manifesto_bans 和 visual_diversity_strategy 到 _content_state.json
    state_file = project_dir / "_content_state.json"
    if state_file.exists():
        with open(state_file, 'r', encoding='utf-8') as f:
            state = json.load(f)
    else:
        state = {}

    state["manifesto_bans"] = manifesto_dict.get("english_cliche_bans", "")
    state["visual_diversity_strategy"] = manifesto_dict.get("visual_diversity_strategy", "")
    with open(state_file, 'w', encoding='utf-8') as f:
        json.dump(state, f, ensure_ascii=False, indent=2)

    # Step 4: 生成人类可审阅的 master_plan.md
    print("\n📋 [Step 4/4] 正在生成完整审阅计划...")
    content_md_path_str = str(content_md_path) if content_md_path.exists() else None
    if content_md_path_str:
        review_md_content = build_master_plan_from_content_plan(
            content_md_path_str, style_config, meta,
            manifesto=manifesto_text,
            per_slide_suggestions=per_slide_suggestions,
        )
    else:
        raise FileNotFoundError(f"content_plan.md 不存在: {content_md_path}。请先运行 plan-content。")
    with open(review_md_path, "w", encoding='utf-8') as f:
        f.write(review_md_content)

    print(f"\n{'='*60}")
    print(f"📋 视觉计划已就绪，已保存: {review_md_path}")
    print(f"   请重点审阅每页的【配图/画面】描述，可直接编辑修改。")
    print(f"{'='*60}")

    print(f"\n{'='*60}")
    print(f"⛔ STOP — 请审阅 {REVIEW_MD_FILENAME} 后再继续！")
    print(f"   必须先让用户审阅：")
    print(f"   ① 视觉主张和配色方案（Design Manifesto）")
    print(f"   ② 每页的配图/画面描述（VisualDirector 已生成，可直接审阅或修改）")
    print(f"{'='*60}")

    # 实际阻塞等待确认（非交互环境返回 "BLOCKED"）
    result = _prompt_user("确认审阅完毕，要继续执行吗？")
    if result == "BLOCKED":
        print(f"""
═══════════════════════════════════════════════════════
⛔ GATE 2 — 已拦截，需人工确认
═══════════════════════════════════════════════════════
   master_plan.md 已生成（第 1 页有完整审阅说明）。

   请将以下内容呈现给用户：
   ① 确认已审阅 master_plan.md
   ② 确认后可回答"继续"我来执行，或回答"preview 1-2"
     我会用 prototype 先跑 1-2 页给您确认风格。
═══════════════════════════════════════════════════════
""")
        return "BLOCKED"
    elif not result:
        print(f"\n⏸️  已停止。请审阅 {REVIEW_MD_FILENAME} 后再运行 execute。")
        return None

    return str(review_md_path)


def generate_plan(content_file: str, template_file: str = None,
                  logo_file: str = None, output_name: str = None, page_count: int = None,
                  style_preference: str = None, briefing: str = None):
    """
    Legacy 入口：与 plan-content 行为一致，仅生成 content_plan.md 并停止。
    请审阅 content_plan.md 确认后，再运行 plan-visual。
    """
    return generate_content_plan(content_file, template_file, logo_file, output_name, page_count, briefing, reuse_existing=False)


# ──────────────────────────────────────────────
#  Phase 2: Execute
# ──────────────────────────────────────────────

def _resolve_execute_input(path_arg: str) -> tuple:
    """解析 execute 输入，支持目录或文件。返回 (plan_path, project_dir, from_review_md)"""
    p = Path(path_arg)
    if not p.exists():
        return None, None, False
    if p.is_dir():
        review_md = p / REVIEW_MD_FILENAME
        plan_json = p / "plan.json"
        # 优先使用 master_plan.md（用户审阅确认的），确保生成内容与大纲一致
        if review_md.exists():
            return str(review_md), str(p), True
        if plan_json.exists():
            return str(plan_json), str(p), False
        print(f"❌ 目录中未找到 {REVIEW_MD_FILENAME} 或 plan.json")
        return None, None, False
    if p.suffix.lower() == ".md":
        name = Path(p).stem
        from datetime import date
        date_prefix = date.today().strftime("%Y%m%d")
        dir_name = f"{date_prefix}_{name}"
        project_dir = DEFAULT_OUTPUT_BASE / dir_name
        project_dir.mkdir(parents=True, exist_ok=True)
        return str(p), str(project_dir), True
    if p.suffix.lower() == ".json":
        return str(p), str(p.parent), False
    print("❌ 请传入 master_plan.md、plan.json 或项目目录")
    return None, None, False


def build_content_only_prototype_pptx(project_dir: str, output_name: str = None) -> str:
    """
    仅根据 master_plan.md 生成纯文字打样 PPT（不生成生图提示词、不调用生图 API）。
    每页包含：标题、副标题（如有）、正文、演讲备注。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    review_md = Path(project_dir) / REVIEW_MD_FILENAME
    if not review_md.exists():
        print(f"❌ 未找到 {REVIEW_MD_FILENAME}，无法生成内容打样")
        return None

    with open(review_md, "r", encoding="utf-8") as f:
        parsed = parse_review_md(f.read())

    pages = parsed.get("pages", [])
    if not pages:
        print("❌ 解析 master_plan 未得到任何页面")
        return None

    name = output_name or Path(project_dir).name
    out_path = Path(project_dir) / f"{name}_prototype_content.pptx"

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        tc = page.get("text_content", {})
        headline = tc.get("headline", "").strip()
        subhead = tc.get("subhead", "").strip()
        body_list = tc.get("body") or []
        speaker_notes = page.get("speaker_notes", "").strip()

        # 标题（顶部）
        if headline:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(15), Inches(1.2))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = headline
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x2b, 0x2b, 0x2b)

        # 副标题（标题下方）
        top = 1.7
        if subhead:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(15), Inches(0.8))
            sf = sub_box.text_frame
            sf.word_wrap = True
            sp = sf.paragraphs[0]
            sp.text = subhead
            sp.font.size = Pt(18)
            sp.font.color.rgb = RGBColor(0x9a, 0x95, 0x93)
            top += 1.0

        # 正文
        body_text = "\n".join(body_list).strip() if body_list else ""
        if body_text:
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(15), Inches(5.5))
            bf = body_box.text_frame
            bf.word_wrap = True
            for i, line in enumerate(body_list):
                if not line.strip():
                    continue
                if i == 0:
                    bp = bf.paragraphs[0]
                else:
                    bp = bf.add_paragraph()
                bp.text = line.replace("**", "").strip()
                bp.font.size = Pt(16)
                bp.font.color.rgb = RGBColor(0x2b, 0x2b, 0x2b)
                bp.space_after = Pt(6)

        # 演讲备注
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    prs.save(str(out_path))
    print(f"✅ 内容打样 PPT 已生成（仅文字，无生图）: {out_path}")
    return str(out_path)


def _auto_select_prototype_seeds(plan_input: str) -> list:
    """Pick one seed page per master family (hero, content, section) for comprehensive prototype preview."""
    from tools.nano_banana_ppt.core.executor import CONTENT_FAMILY, BOOKEND_FAMILY

    plan_path, proj_dir, from_review = _resolve_execute_input(plan_input)
    if not plan_path:
        return [1, 2]

    pages = []
    if from_review:
        with open(plan_path, "r", encoding="utf-8") as f:
            parsed = parse_review_md(f.read())
        pages = parsed.get("pages", [])
    else:
        with open(plan_path, "r", encoding="utf-8") as f:
            raw = json.load(f)
        pages = raw if isinstance(raw, list) else raw.get("slides", [])

    if not pages:
        return [1, 2]

    FAMILY_MAP = {
        'hero': 'hero',
        'section': 'section',
    }
    for t in BOOKEND_FAMILY:
        FAMILY_MAP[t] = 'bookend'
    for t in CONTENT_FAMILY:
        FAMILY_MAP[t] = 'content'

    selected = {}
    for p in pages:
        p_type = p.get('type', 'content').lower()
        if p_type.startswith('template_') or p_type == 'background_only':
            continue
        family = FAMILY_MAP.get(p_type, 'content')
        if family not in selected:
            selected[family] = p.get('page_num', 1)

    seeds = sorted(selected.values())
    return seeds if seeds else [1, 2]


def execute_from_plan(plan_input: str, output_name: str = None, resolution: str = "1K", slide_filter: list = None, reassemble_only: bool = False, regenerate: bool = False, no_blend: bool = False, force: bool = False):
    """
    Phase 2: 执行阶段。
    接受 REVIEW_MD_FILENAME 或 plan.json 路径，或包含它们的目录，进行最终的图文生成与组装。

    plan.json 复用策略：
    - 若 plan.json 已存在且未指定 --regenerate，则直接复用（无论全量还是补跑）
    - 若 plan.json 不存在，或指定了 --regenerate，则从 master_plan.md 重新派生
    - 派生现在是模板拼接方式（无 LLM 调用，即时完成）
    """
    # 获取 API 配置
    api_key = os.getenv("OPENAI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    api_base = os.getenv("OPENAI_API_BASE")

    plan_path, proj_dir, from_review = _resolve_execute_input(plan_input)
    if not plan_path:
        return None

    # GATE 2: 在非交互模式下，如果没有通过确认流程，则阻止直接执行
    if from_review and not sys.stdin.isatty() and not force:
        print("""
⚠️  执行保护：检测到您直接运行 execute 命令（未经过审阅确认流程）。

   master_plan.md 已生成，请确认：
   ① 内容是否已审阅？
   ② 是否需要用 prototype 预览 1-2 页后再全量生成？

   如需强制执行，请加 --force 参数：
   python main.py execute \"{proj_dir}\" --force
""".format(proj_dir=proj_dir))
        return None

    plan_json_path = Path(proj_dir) / "plan.json"

    # 问题4: 内容一致性校验
    if not force:
        issues = validate_content_visual_consistency(Path(proj_dir))
        if issues:
            print("\n❌ 发现内容不一致问题:")
            for issue in issues:
                print(f"   - {issue}")
            print("\n请先修复这些问题，或使用 --force 跳过校验")
            return None
        else:
            print("\n✅ 内容一致性校验通过")

    # 智能检测文件修改时间，自动决定是否需要重新生成
    use_existing_plan = False
    if plan_json_path.exists() and not regenerate:
        if from_review:
            md_mtime = Path(plan_path).stat().st_mtime
            json_mtime = plan_json_path.stat().st_mtime
            if json_mtime > md_mtime:
                use_existing_plan = True
                print("\n📄 复用已有 plan.json（比 master_plan.md 更新）")
            else:
                print("\n📄 检测到 master_plan.md 已更新，重新生成 plan.json")
        else:
            use_existing_plan = True
            print("\n📄 复用已有 plan.json（如需重新生成提示词，请加 --regenerate 参数）")

    if from_review:
        if use_existing_plan:
            plan_path = str(plan_json_path)
        else:
            reset_session()
            print("\n📄 正在从 master_plan.md 派生技术计划（模板拼接，无 LLM 调用）...")
            with open(plan_path, "r", encoding="utf-8") as f:
                md_text = f.read()
            parsed = parse_review_md(md_text)
            if not parsed.get("pages"):
                print("❌ 解析失败，未找到有效页面")
                return None
            pmeta = parsed.get("meta", {})
            content_file = pmeta.get("content_file", "")
            plan_data = derive_technical_plan(parsed, proj_dir, content_file, api_key, api_base)
            with open(plan_json_path, "w", encoding="utf-8") as f:
                json.dump(plan_data, f, ensure_ascii=False, indent=2)
            print(f"✅ 技术计划已保存: {plan_json_path}")
            plan_path = str(plan_json_path)
    else:
        # 直接传入 plan.json 时：如果 master_plan.md 也存在且 plan.json 需要重新生成
        review_md = Path(proj_dir) / REVIEW_MD_FILENAME
        if review_md.exists() and not use_existing_plan:
            print("\n📄 正在从 master_plan.md 重新派生技术计划...")
            with open(review_md, "r", encoding="utf-8") as f:
                md_text = f.read()
            parsed = parse_review_md(md_text)
            if parsed.get("pages"):
                pmeta = parsed.get("meta", {})
                content_file = pmeta.get("content_file", "")
                plan_data = derive_technical_plan(parsed, proj_dir, content_file, api_key, api_base)
                with open(plan_json_path, "w", encoding="utf-8") as f:
                    json.dump(plan_data, f, ensure_ascii=False, indent=2)
                print(f"✅ 技术计划已更新并保存: {plan_json_path}")
                plan_path = str(plan_json_path)

    with open(plan_path, 'r', encoding='utf-8') as f:
        plan_data = json.load(f)

    if isinstance(plan_data, list):
        slides = plan_data
        meta = {}
    else:
        slides = plan_data.get("slides", plan_data)
        meta = plan_data.get("meta", {})

    name = output_name or meta.get("project_name") or Path(meta.get("content_file", "presentation")).stem
    project_dir = meta.get("project_dir") or proj_dir
    if not project_dir:
        from datetime import date
        date_prefix = date.today().strftime("%Y%m%d")
        dir_name = f"{date_prefix}_{name}"
        project_dir = str(DEFAULT_OUTPUT_BASE / dir_name)
    
    template_path = meta.get("template_file")

    print(f"\n⚡ Nano Banana 2 — Phase 2: Execute")
    print(f"📋 计划文件: {plan_path} ({len(slides)} 页)")
    print(f"📂 项目目录: {project_dir}")
    print(f"📦 输出名称: {name}")
    print(f"🖼️ 分辨率: {resolution or '1K'}")

    # 写一份纯 slides list 供 executor 读取
    exec_plan_file = str(Path(project_dir) / "_exec_slides.json")
    Path(project_dir).mkdir(parents=True, exist_ok=True)
    with open(exec_plan_file, "w", encoding='utf-8') as f:
        json.dump(slides, f, ensure_ascii=False, indent=2)

    out = execute_plan(
        exec_plan_file, name, template_path, project_dir=project_dir,
        resolution=resolution or "1K", slide_filter=slide_filter,
        reassemble_only=reassemble_only
    )

    if os.path.exists(exec_plan_file):
        os.remove(exec_plan_file)

    return out, project_dir


def _parse_slides_arg(s: str) -> list:
    """解析用户输入的页号，如 '3 5 7' 或 '3,5,7'"""
    if not s or not s.strip():
        return None
    nums = re.findall(r'\d+', s)
    return [int(n) for n in nums] if nums else None


def _prompt_user(message: str, default_yes: bool = True) -> bool | str:
    """
    向用户提示确认。在非交互终端（AI Agent / Claude Code）下返回 "BLOCKED"。
    """
    # 非交互环境：返回特殊值，让上游知道是被拦截而非用户拒绝
    if not sys.stdin.isatty():
        print(f"\n[非交互模式: {message} — 已拦截，需用户介入]")
        return "BLOCKED"
    hint = "(Y/n)" if default_yes else "(y/N)"
    try:
        answer = input(f"\n{message} {hint}: ").strip().lower()
    except EOFError:
        return False
    if answer == '':
        return default_yes
    return answer in ('y', 'yes')


def _interactive_rerun_prompt(plan_input: str, output_name: str, resolution: str, project_dir: str = None):
    """生成完成后，交互式询问是否需要重跑部分页面。rerun 使用 project_dir 以命中 plan.json"""
    print("\n" + "=" * 50)
    print("是否需要重跑部分页面？")
    print("  输入页号（如 3 5 7 或 3,5,7），直接回车跳过")
    print("=" * 50)
    try:
        user_input = input("> ").strip()
    except EOFError:
        return
    slides = _parse_slides_arg(user_input)
    if not slides:
        return
    rerun_input = project_dir or plan_input
    execute_from_plan(rerun_input, output_name, resolution=resolution, slide_filter=slides)
    _interactive_rerun_prompt(plan_input, output_name, resolution, project_dir=project_dir)


# ──────────────────────────────────────────────
#  Legacy: 一键全自动
# ──────────────────────────────────────────────

def auto_generate_ppt(content_file: str, template_file: str = None,
                      logo_file: str = None, output_name: str = None, resolution: str = "1K", page_count: int = None,
                      style_preference: str = None, briefing: str = None):
    """
    交互式终端下的一键全自动流程。
    在每个阶段完成后必须等待用户确认才继续。
    非交互环境（如 Cursor Shell/Agent）下自动停止，提示用户手动执行下一步。
    """
    # ── Phase 1.1: 内容规划 ──
    content_md = generate_content_plan(content_file, template_file, logo_file, output_name, page_count=page_count, briefing=briefing, reuse_existing=False)
    if not content_md:
        return

    project_dir = str(Path(content_md).parent)

    if not _prompt_user("请审阅 content_plan.md，确认内容无误后继续生成视觉计划？"):
        print(f"\n⏸️  已停止。请审阅后手动执行下一步：")
        print(f"   python -m tools.nano_banana_ppt.main plan-visual \"{project_dir}\" [--style xxx]")
        return

    # ── Phase 1.5: 视觉规划 ──
    plan_file = generate_visual_plan(project_dir, style_preference=style_preference)
    if not plan_file:
        return

    if not _prompt_user("请审阅 master_plan.md，确认制作蓝图无误后继续执行生图？"):
        print(f"\n⏸️  已停止。请审阅后手动执行下一步：")
        print(f"   python -m tools.nano_banana_ppt.main execute \"{project_dir}\"")
        return

    # ── Phase 2: 执行 ──
    resolution = (resolution or "1K").strip().upper()
    if resolution not in ("1K", "2K", "4K"):
        resolution = "1K"

    out_name = output_name or Path(plan_file).parent.name
    result = execute_from_plan(plan_file, out_name, resolution=resolution)
    out_path, proj_dir = result if isinstance(result, tuple) else (result, None)
    _interactive_rerun_prompt(plan_file, out_name, resolution, project_dir=proj_dir)


def execute_upscale(proj_dir: str, resolution: str, slide_filter: list = None):
    """独立的高保真放大流程"""
    if not os.path.exists(proj_dir):
        print(f"❌ 项目目录不存在: {proj_dir}")
        return False

    resolution = (resolution or "4K").upper()
    if resolution not in ("2K", "4K"):
        print(f"❌ Upscale 分辨率参数错误 ({resolution})，仅支持 2K 或 4K。")
        return False

    print(f"\n🔍 开始执行高保真 Upscale 放大流程...")
    print(f"📁 项目目录: {proj_dir}")
    print(f"🖼️ 目标分辨率: {resolution}")
    if slide_filter:
        print(f"📑 仅放大指定页面: {slide_filter}")
        
    api_key = os.getenv("OPENAI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    api_base = os.getenv("OPENAI_API_BASE")
    if not api_key:
        print("❌ 未设置 OPENAI_API_KEY 或 GOOGLE_API_KEY")
        return False

    from tools.nano_banana_ppt.core.generator import PPTGenerator
    generator = PPTGenerator(api_key=api_key, api_base=api_base, slides_dir=proj_dir)
    
    slides_dir = Path(proj_dir) / "slides"
    if not slides_dir.exists():
        # 兼容旧版本可能直接生成在根目录的情况
        slides_dir = Path(proj_dir)
        
    # 查找背景图片
    bg_files = list(slides_dir.glob("slide_*.png")) + list(slides_dir.glob("slide_*.jpg"))
    if not bg_files:
        print(f"❌ 在项目目录中未找到任何图片 (slide_*.png/jpg)")
        return False
        
    # 按照页码排序
    bg_files.sort()
    
    success_count = 0
    total_count = 0
    
    for bg_file in bg_files:
        # 从文件名提取页码 (如 slide_01.png -> 1)
        try:
            filename = bg_file.stem
            page_num_str = filename.replace("slide_", "")
            page_num = int(page_num_str)
        except ValueError:
            continue
            
        # 过滤页码
        if slide_filter and page_num not in slide_filter:
            continue
            
        total_count += 1
        print(f"[{page_num}] 正在放大: {bg_file.name} ...")
        
        success = generator.upscale_image(str(bg_file), resolution=resolution)
        if success:
            success_count += 1
            
    print(f"\n✅ Upscale 图片放大完成! 成功/总数: {success_count}/{total_count}")
    
    if success_count > 0:
        print("\n🔨 正在使用高清图片重新组装 PPTX...")
        output_name = Path(proj_dir).name
        # 复用已有的 execute_from_plan 的组装逻辑
        # 传递 proj_dir 作为 plan_input，设置 reassemble_only=True
        execute_from_plan(proj_dir, output_name, resolution=resolution, slide_filter=slide_filter, reassemble_only=True)
        print("🎉 高清 PPTX 组装完成!")
        
    return True

# ──────────────────────────────────────────────
#  CLI 入口
# ──────────────────────────────────────────────

def print_usage():
    print("""
Nano Banana 2 PPT Generator

用法:
  # Phase 1: 内容规划（生成 content_plan.md 供大纲审阅）
  python -m tools.nano_banana_ppt.main plan-content <content_file> [template_file] [logo_file] [output_name] [--pages N]

  # Phase 1.5: 制作蓝图（生成 master_plan.md 包含内容+风格+视觉主张）
  python -m tools.nano_banana_ppt.main plan-visual <项目目录> [--style "风格描述或预设名称"]

  # Prototype: 快速原型验证（仅生成前两页或指定页）；加 --content-only 则仅生成纯文字打样（不生图、不生成生图提示词）
  python -m tools.nano_banana_ppt.main prototype <项目目录或plan文件> [output_name] [--slides 1 2] [--resolution 1K|2K|4K] [--content-only]

  # Phase 2: 执行计划（生图 + 组装 PPTX）
  # 可传入项目目录、master_plan.md 或 plan.json
  python -m tools.nano_banana_ppt.main execute <项目目录或plan文件> [output_name] [--resolution 1K|2K|4K] [--slides 3 5 7] [--reassemble] [--force]

  # Upscale: 后置高保真放大 (1K -> 2K/4K)
  python -m tools.nano_banana_ppt.main upscale <项目目录> [--resolution 2K|4K] [--slides 3 5 7]

  # 一键全自动（交互式终端下可用）
  python -m tools.nano_banana_ppt.main auto <content_file> [template_file] [logo_file] [output_name] [--resolution 1K|2K|4K]

参数:
  --pages       期望页数，如 --pages 15
  --style       风格偏好，如 --style "毕加索立体主义"
  --briefing    用户意图，如 --briefing "我最想传达的是：投资需警惕左尾风险"
  --resolution  分辨率，1K|2K|4K，默认 1K
  --slides      仅重跑指定页号，如 --slides 3 5 7
  --reassemble  仅从已有 slides 重新组装 PPTX，不生成新图（用于修复布局问题）
  --regenerate  强制重新从 master_plan.md 派生 plan.json（覆盖已有缓存）
  --no-blend    跳过原生图片的 Blend Pass（打样阶段加速用，不影响最终质量）
  --force       跳过执行保护确认（在非交互模式下绕过 GATE 2 警告）

交互:
  执行完成后（交互式终端），可输入页号（如 3 5 7）重跑部分页面，直接回车跳过。

目录结构:
  ~/Desktop/AI output/ppt/{date}_{name}/
    {date}_{name}.pptx               ← 最终交付的 PPTX
    content_plan.md                  ← 内容草稿（Phase 1.1 生成）
    master_plan.md                   ← 制作蓝图（Phase 1.5 生成，含内容+视觉）
    plan.json                        ← 技术计划
    template_assets/                 ← 模版提取物
    slides/                          ← 生成的页面图片
""")


def _parse_cli_args(args):
    """解析 CLI 参数，提取 --resolution、--slides、--pages、--style、--briefing、--reassemble、--content-only、--regenerate、--no-blend、--force"""
    rest = []
    resolution = "1K"  # 默认值为 1K
    slides = None
    pages = None
    style = None
    briefing = None
    reassemble = False
    content_only = False
    regenerate = False
    no_blend = False
    force = False
    reuse_existing = False
    i = 0
    while i < len(args):
        a = args[i]
        if a == "--pages" and i + 1 < len(args):
            pages = int(args[i + 1])
            i += 2
        elif a == "--style" and i + 1 < len(args):
            style = args[i + 1]
            i += 2
        elif a == "--briefing" and i + 1 < len(args):
            briefing = args[i + 1]
            i += 2
        elif a == "--resolution" and i + 1 < len(args):
            resolution = args[i + 1].upper()  # 统一转换为大写
            if resolution not in ("1K", "2K", "4K"):
                resolution = "1K"  # 无效值时使用默认值
            i += 2
        elif a == "--reassemble":
            reassemble = True
            i += 1
        elif a == "--content-only":
            content_only = True
            i += 1
        elif a == "--regenerate":
            regenerate = True
            i += 1
        elif a == "--no-blend":
            no_blend = True
            i += 1
        elif a == "--force":
            force = True
            i += 1
        elif a == "--reuse-existing":
            reuse_existing = True
            i += 1
        elif a == "--slides":
            i += 1
            nums = []
            while i < len(args):
                if args[i].isdigit():
                    nums.append(int(args[i]))
                    i += 1
                elif ',' in args[i] and all(p.strip().isdigit() for p in args[i].split(',') if p.strip()):
                    nums.extend([int(p.strip()) for p in args[i].split(',') if p.strip()])
                    i += 1
                else:
                    break
            slides = nums if nums else None
        else:
            rest.append(a)
            i += 1
    return rest, resolution, slides, pages, style, briefing, reassemble, content_only, regenerate, no_blend, force, reuse_existing


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print_usage()
        sys.exit(1)

    rest, resolution, slides, pages, style, briefing, reassemble, content_only, regenerate, no_blend, force, reuse_existing = _parse_cli_args(sys.argv[1:])
    command = rest[0].lower() if rest else ""

    if command == "plan" or command == "plan-content":
        if len(rest) < 2:
            print("❌ 缺少 content_file 参数")
            sys.exit(1)
        content = rest[1]
        template = rest[2] if len(rest) > 2 else None
        logo = rest[3] if len(rest) > 3 else None
        out_name = rest[4] if len(rest) > 4 else None
        if command == "plan":
            generate_plan(content, template, logo, out_name, page_count=pages, briefing=briefing)
        else:
            generate_content_plan(content, template, logo, out_name, page_count=pages, briefing=briefing, reuse_existing=reuse_existing)

    elif command == "plan-visual":
        if len(rest) < 2:
            print("❌ 缺少 项目目录 参数")
            sys.exit(1)
        proj_dir = rest[1]
        result = generate_visual_plan(proj_dir, style_preference=style)
        if result == "BLOCKED":
            # 已拦截，Agent 应停止并向用户呈现 plan 等待确认
            sys.exit(0)  # 干净退出，不抛异常

    elif command == "execute":
        if len(rest) < 2:
            print("❌ 缺少 plan_file 或项目目录参数")
            sys.exit(1)
        pf = rest[1]
        on = rest[2] if len(rest) > 2 else None
        result = execute_from_plan(pf, on, resolution=resolution, slide_filter=slides, reassemble_only=reassemble, regenerate=regenerate, no_blend=no_blend, force=force)
        if result:
            _, proj_dir = result if isinstance(result, tuple) else (None, None)
            _interactive_rerun_prompt(pf, on or Path(pf).parent.name, resolution or "1K", project_dir=proj_dir)

    elif command == "upscale":
        if len(rest) < 2:
            print("❌ 缺少项目目录参数")
            sys.exit(1)
        proj_dir = rest[1]
        
        # 验证输入确实是个目录
        if not os.path.isdir(proj_dir):
            print(f"❌ {proj_dir} 不是一个有效的目录。upscale 需要传入包含图片的 ~/Desktop/AI output/ppt/项目目录。")
            sys.exit(1)
            
        # 默认为 4K 放大
        target_res = resolution if resolution in ("2K", "4K") else "4K"
        execute_upscale(proj_dir, resolution=target_res, slide_filter=slides)

    elif command == "auto":
        content = rest[1] if len(rest) > 1 else None
        if not content:
            print("❌ 缺少 content_file 参数")
            sys.exit(1)
        template = rest[2] if len(rest) > 2 else None
        logo = rest[3] if len(rest) > 3 else None
        output = rest[4] if len(rest) > 4 else None
        auto_generate_ppt(content, template, logo, output, resolution=resolution, page_count=pages, style_preference=style, briefing=briefing)

    elif command == "prototype":
        if len(rest) < 2:
            print("❌ 缺少项目目录或plan文件参数")
            sys.exit(1)
        pf = rest[1]
        output_name = rest[2] if len(rest) > 2 else None

        if content_only:
            # 仅生成打样 PPT 内容（纯文字），不生成生图提示词、不调用生图 API
            plan_path, proj_dir, _ = _resolve_execute_input(pf)
            if not plan_path or not proj_dir:
                print("❌ 无法解析项目目录或 master_plan.md")
                sys.exit(1)
            print("\n🚀 Nano Banana 2 — 内容打样（仅文字，不生图）")
            build_content_only_prototype_pptx(proj_dir, output_name or Path(proj_dir).name)
        else:
            if slides is not None:
                target_slides = slides
            else:
                target_slides = _auto_select_prototype_seeds(pf)
            print(f"\n🚀 Nano Banana 2 — Prototype Mode")
            print(f"   自动选取每种母版类型的种子页进行视觉验证: {target_slides}")
            result = execute_from_plan(pf, output_name, resolution=resolution, slide_filter=target_slides, regenerate=regenerate, no_blend=no_blend, force=True)
            if result:
                print("\n✨ Prototype complete. Check the slides directory. If you are satisfied with the style, run the 'execute' command to generate the full presentation.")

    else:
        print(f"❌ 未知命令: {command}")
        print_usage()
        sys.exit(1)
